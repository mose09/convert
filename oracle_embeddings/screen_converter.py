"""화면변환기 PoC — AS-IS 화면 캡처 → TO-BE PPTX (Vision LLM).

기존 자산 재사용:
 - ``legacy_pattern_discovery._call_llm(image_paths=...)`` : OpenAI 호환
   multimodal vision 호출 (base64 인코딩 + JSON 파싱 + 재시도 + raw dump
   인프라). 본 모듈은 system_prompt 만 주입해서 재사용.
 - ``python-pptx`` : TO-BE 슬라이드 도형/표/텍스트박스 렌더.

설계 원칙 (PoC):
 - 캐시 없음. 매번 VLM 호출.
 - 템플릿은 layout 추출용 비주얼 컨텍스트로만 사용 (별도 스타일 스펙
   추출 단계 없음). VLM 이 AS-IS 와 템플릿 캡처들을 동시에 보고 layout
   JSON 만 뽑는다.
 - 도형 스타일 고정 (검정 선/흰 배경/맑은 고딕). 색상/폰트의 템플릿
   반영은 PoC 동작 확인 후 후속.

모델 선택:
 - ``_call_llm`` 우선순위 (``PATTERN_LLM_MODEL`` > ``LLM_MODEL`` >
   config.yaml ``llm.model``) 를 그대로 따른다. 화면변환 전용 env 키는
   두지 않고 기존 키를 공유한다. vision 가능 모델을 쓰려면
   ``PATTERN_LLM_MODEL`` 또는 ``LLM_MODEL`` 을 Qwen3-VL/2.5-VL 등으로
   지정해야 한다.
"""
from __future__ import annotations

import json
import logging
import os
import re
from pathlib import Path
from typing import Any

from .legacy_pattern_discovery import _call_llm

logger = logging.getLogger(__name__)

_IMG_EXTS = (".png", ".jpg", ".jpeg")

_SYSTEM_PROMPT = (
    "당신은 레거시 화면 캡처 분석가입니다. 사용자가 제공한 첫 번째 이미지는 "
    "AS-IS 화면 캡처이고, 그 뒤의 이미지들은 TO-BE 디자인 템플릿 캡처입니다. "
    "AS-IS 의 비즈니스 의미는 유지하면서 템플릿의 시각 스타일을 따르는 "
    "TO-BE 레이아웃을 JSON 으로만 출력하세요. 다른 설명/마크다운/코멘트 금지."
)

_USER_PROMPT = """다음 스키마로 TO-BE 화면 레이아웃 JSON 을 출력해줘. 키:

{
  "page_title": "화면 타이틀 한 줄",
  "search_fields": [{"label": "조건명", "type": "text|select|date|checkbox"}],
  "table_columns": ["컬럼1", "컬럼2"],
  "buttons": ["조회", "초기화", "저장"],
  "notes": "특이사항 1~2줄",
  "regions": {
    "title":         {"x": "<0~1>", "y": "<0~1>", "w": "<0~1>", "h": "<0~1>"},
    "search_panel":  {"x": "<0~1>", "y": "<0~1>", "w": "<0~1>", "h": "<0~1>", "cols": "<int>"},
    "table":         {"x": "<0~1>", "y": "<0~1>", "w": "<0~1>", "h": "<0~1>"},
    "buttons":       {"x": "<0~1>", "y": "<0~1>", "w": "<0~1>", "h": "<0~1>", "align": "left|center|right"},
    "notes":         {"x": "<0~1>", "y": "<0~1>", "w": "<0~1>", "h": "<0~1>"}
  }
}

규칙:
- AS-IS 에 없는 필드/버튼/컬럼은 만들지 마세요 (할루시네이션 금지).
- 모든 라벨은 한국어 원문 유지.
- 비어 있는 키는 빈 배열 [] / 빈 문자열 "" 로.
- JSON 객체 하나만 출력.
- **React 소스가 첨부된 경우 (프롬프트 하단에 `=== React 소스` 블록)**
  → 그 소스가 `search_fields` / `table_columns` / `buttons` 의 **정답**
  (JSX 의 `<input label=...>` / `<Column header=...>` / `<Button>...`
  텍스트 그대로 옮길 것). 이미지에서 다르게 보여도 소스를 우선.
  이미지는 `regions` (위치/크기) 와 `page_title` 추론에만 사용.
- 소스가 첨부되지 않으면 AS-IS 이미지에서 모든 항목 추출.

**순서 규칙 — 매우 중요**:
- `search_fields` 는 화면을 사람이 읽는 자연 순서 = **row-major
  (행 우선)** 로 나열할 것. 즉 같은 행에 있는 필드들을 **왼쪽 → 오른쪽**
  순서로 먼저 나열한 뒤, 다음 행으로 내려가서 같은 방식으로 나열.
  **세로(컬럼) 단위로 위→아래 먼저 훑는 column-major 순서는 절대 금지.**
- `table_columns` 는 표 헤더의 **왼쪽 → 오른쪽** 순서 그대로.
- `buttons` 는 화면에 보이는 **왼쪽 → 오른쪽** (같은 행이면)
  또는 위→아래 순서로 나열.

예시 — 검색 패널이 3열 × 2행 그리드로 다음과 같이 배치되어 있다면:

  [주문번호] [주문일자] [상태]
  [고객사  ] [품목    ] [금액]

올바른 출력 (row-major):
  ["주문번호", "주문일자", "상태", "고객사", "품목", "금액"]

잘못된 출력 (column-major, 사용 금지):
  ["주문번호", "고객사", "주문일자", "품목", "상태", "금액"]

**regions 규칙 — 위치/크기 정확도의 핵심**:
- 모든 좌표는 슬라이드 전체를 1.0 으로 한 **normalized 값 (0.0 ~ 1.0)**.
  top-left origin (x=0 가 왼쪽, y=0 가 위).
- `x`, `y` = 영역의 좌상단 모서리. `w`, `h` = 영역의 가로/세로 크기.
  스키마의 `<0~1>` 은 자리 표시일 뿐 — **반드시 템플릿 캡처를 직접 보고
  실제 관찰된 좌표값(float)을 넣어라**. 스키마 placeholder 를 그대로
  복사하면 안 됨.
- bbox 의 기준은 **TEMPLATE 캡처의 영역 위치/크기** (TO-BE 디자인).
  AS-IS 의 위치는 무시하고, 템플릿이 보여주는 레이아웃 비율을 따르세요.
- `search_panel.cols` = 패널 안 필드 그리드의 한 행에 들어가는 컬럼 수
  (보통 3 ~ 6). 필드를 몇 줄로 배치할지 결정.
- `buttons.align` = `"left"` | `"center"` | `"right"`. 템플릿 보고 결정.

**buttons.y 결정 규칙 — 매우 중요 (가장 많이 틀리는 부분)**:
- 버튼의 y 는 **템플릿에서 버튼 행이 실제로 위치한 자리** 를 따라가야
  한다. 화면 종류마다 다르다:
  * 조회/필터 버튼이 **검색 패널 위 또는 같은 행 오른쪽** 에 있으면
    → `buttons.y` 는 0.05 ~ 0.20 (작은 값, 화면 윗부분)
  * 조회 버튼이 **검색 패널 안 마지막 칸** 에 있으면
    → search_panel 의 y 범위 안에 위치 (보통 0.20 ~ 0.30)
  * 저장/삭제 등이 **표 아래** 에 있으면 → 0.80 ~ 0.95 (큰 값, 아랫부분)
- **절대 기본값 0.86 등 큰 y 를 자동으로 쓰지 마라.** 템플릿을 직접 보고
  실제 버튼 행의 세로 위치를 측정해 넣어라.

- 패널/테이블/버튼바가 겹치지 않도록 y 좌표 분리. 화면에 없는 영역은
  생략 가능 (예: 표가 없는 입력 폼은 `table` 키 자체를 빼도 됨)."""


# ── 입력 수집 ────────────────────────────────────────────────────────


def _list_images(folder: Path) -> list[Path]:
    if not folder.is_dir():
        raise SystemExit(f"폴더 없음 또는 디렉토리 아님: {folder}")
    return sorted(p for p in folder.iterdir()
                  if p.is_file() and p.suffix.lower() in _IMG_EXTS)


# ── 스타일 프로파일 ──────────────────────────────────────────────────

# 템플릿 미제공 또는 추출 실패 시 사용되는 기본값 — 기존 하드코드와 동일.
_DEFAULT_STYLE: dict[str, str] = {
    "primary_color": "#1F3A5F",          # 헤더/주요 강조
    "title_color": "#1F3A5F",
    "section_label_color": "#444444",
    "panel_bg": "#F5F7FA",
    "panel_border": "#CFD6E0",
    "input_bg": "#FFFFFF",
    "input_border": "#A0A8B4",
    "field_label_color": "#333333",
    "input_placeholder_color": "#999999",
    "table_header_bg": "#1F3A5F",
    "table_header_text": "#FFFFFF",
    "table_row_bg": "#FFFFFF",
    "button_bg": "#1F3A5F",
    "button_text": "#FFFFFF",
    "button_shape": "rounded",           # rounded | square
    "font_family": "맑은 고딕",
    "notes_color": "#555555",
}

_STYLE_SYSTEM_PROMPT = (
    "당신은 UI 디자인 스타일 추출기입니다. 첨부된 디자인 템플릿 캡처 "
    "이미지들을 보고 그 시각 스타일 (색·폰트·버튼 모양) 을 JSON 으로만 "
    "출력합니다. 설명/마크다운/코멘트 없이 JSON 객체 하나만."
)

_STYLE_PROMPT = """첨부된 템플릿 캡처들의 시각 스타일을 다음 스키마로 추출.
모든 화면에 일관 적용할 거니까 가장 빈번한/대표적인 값을 선택.

스키마 (모든 색은 "#RRGGBB" 6자리 hex):
{
  "primary_color":          "#RRGGBB",   // 강조 (헤더/주요 버튼)
  "title_color":            "#RRGGBB",   // 페이지 타이틀 글자
  "section_label_color":    "#RRGGBB",   // "검색 조건" 등 섹션 라벨
  "panel_bg":               "#RRGGBB",   // 검색/입력 패널 배경
  "panel_border":           "#RRGGBB",   // 패널 보더
  "input_bg":               "#RRGGBB",   // 입력 박스 배경
  "input_border":           "#RRGGBB",   // 입력 박스 보더
  "field_label_color":      "#RRGGBB",   // 필드 라벨 글자
  "input_placeholder_color":"#RRGGBB",   // placeholder/타입 힌트 글자
  "table_header_bg":        "#RRGGBB",   // 표 헤더 배경
  "table_header_text":      "#RRGGBB",   // 표 헤더 글자
  "table_row_bg":           "#RRGGBB",   // 표 본문 행 배경
  "button_bg":              "#RRGGBB",   // 버튼 배경
  "button_text":            "#RRGGBB",   // 버튼 글자
  "button_shape":           "rounded|square",  // 버튼 모서리
  "font_family":            "<폰트명>",  // 본문 폰트 (한글 화면이면 맑은 고딕/나눔고딕 등)
  "notes_color":            "#RRGGBB"    // 노트/캡션 글자
}

규칙:
- 정확한 hex 색을 자신할 수 없으면 가장 가까운 추정치라도 출력.
- 모르는 키는 생략 가능 (생략된 키는 기본값 사용).
- JSON 객체 하나만, 설명 없이."""


def _parse_hex_color(s, default_rgb: tuple[int, int, int]):
    """`#RRGGBB` → RGBColor. 파싱 실패 시 default_rgb 로 fallback."""
    from pptx.dml.color import RGBColor
    if isinstance(s, str):
        cleaned = s.strip().lstrip("#")
        if len(cleaned) == 6:
            try:
                r = int(cleaned[0:2], 16)
                g = int(cleaned[2:4], 16)
                b = int(cleaned[4:6], 16)
                return RGBColor(r, g, b)
            except ValueError:
                pass
    return RGBColor(*default_rgb)


def _hex_to_rgb_tuple(s: str) -> tuple[int, int, int]:
    """Default RGB 추출용 헬퍼 (테스트용)."""
    cleaned = s.lstrip("#")
    return (int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


def _resolve_style(style: dict | None) -> dict:
    """사용자 style 과 기본값 병합. 빈/None 이면 전부 기본값."""
    merged = dict(_DEFAULT_STYLE)
    if isinstance(style, dict):
        for k, v in style.items():
            if v is None:
                continue
            if isinstance(v, str) and not v.strip():
                continue
            merged[k] = v
    return merged


def extract_style_profile(template_images: list[Path], config: dict) -> dict:
    """템플릿 캡처들로부터 1회 VLM 호출 → style dict. 실패 시 빈 dict.

    convert() 시작 시 한 번만 호출. 화면별 추출과 분리되어 모든 슬라이드에
    일관 적용된다.
    """
    if not template_images:
        return {}
    print("  스타일 프로파일 추출 (템플릿 1회 분석)...")
    result = _call_llm(
        prompt=_STYLE_PROMPT,
        config=config,
        label="screen_style_profile",
        system_prompt=_STYLE_SYSTEM_PROMPT,
        image_paths=[str(t) for t in template_images],
    )
    if not isinstance(result, dict):
        logger.warning("style 응답이 dict 가 아님 — 기본 스타일 사용")
        return {}
    return result


# ── CSS 스타일 파싱 (TO-BE 스타일 가이드 파일 — LLM 비용 0) ──────────────


_CSS_VAR_ALIASES: dict[str, list[str]] = {
    "primary_color": ["primary-color", "color-primary", "brand-primary",
                       "brand-color", "main-color", "color-main", "color-brand",
                       "primary", "theme-primary"],
    "title_color": ["title-color", "color-title", "heading-color",
                     "color-heading", "text-heading", "h1-color"],
    "section_label_color": ["section-color", "section-label-color",
                              "label-color", "color-label", "secondary-text"],
    "panel_bg": ["panel-bg", "panel-background", "search-bg", "filter-bg",
                  "form-bg", "card-bg", "color-bg-panel"],
    "panel_border": ["panel-border", "border-color", "color-border"],
    "input_bg": ["input-bg", "field-bg", "form-field-bg", "color-bg-input"],
    "input_border": ["input-border", "field-border", "form-field-border"],
    "field_label_color": ["field-label-color", "input-label-color"],
    "input_placeholder_color": ["placeholder-color", "color-placeholder"],
    "table_header_bg": ["table-header-bg", "thead-bg", "grid-header-bg",
                          "header-bg"],
    "table_header_text": ["table-header-text", "thead-text",
                            "grid-header-text", "header-text"],
    "table_row_bg": ["table-row-bg", "tbody-bg", "grid-row-bg", "row-bg"],
    "button_bg": ["button-bg", "btn-bg", "primary-button-bg",
                   "color-btn-primary-bg"],
    "button_text": ["button-text", "btn-text", "primary-button-text",
                     "btn-color", "color-btn-primary-text"],
    "font_family": ["font-family-base", "font-family", "font-sans",
                     "font-default", "base-font"],
    "notes_color": ["notes-color", "caption-color", "color-caption"],
}


def parse_css_style(css_path: Path | str) -> dict:
    """CSS 파일 → style_profile dict.

    추출 우선순위:
      1. CSS 변수 (``:root { --foo: #fff }``) — alias 군으로 first-match
      2. 의미 있는 클래스 (``.btn-primary``, ``thead`` 등) 의 속성
      3. ``body`` / ``:root`` 의 ``font-family`` (첫 폰트만)

    매칭 안 된 키는 dict 에 미포함 → ``_resolve_style`` 이 default 로 fallback.
    LLM 호출 0, 결정적, 폐쇄망 친화.
    """
    p = Path(css_path)
    if not p.is_file():
        return {}
    try:
        text = p.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        logger.warning("CSS 읽기 실패 %s: %s", p, e)
        return {}

    # 주석 제거
    text = re.sub(r"/\*.*?\*/", "", text, flags=re.DOTALL)

    out: dict[str, str] = {}

    # 1. CSS 변수 인덱스
    variables: dict[str, str] = {}
    for m in re.finditer(r"--([\w-]+)\s*:\s*([^;]+);", text):
        variables[m.group(1).strip().lower()] = m.group(2).strip()

    def _resolve_value(v: str, _depth: int = 0) -> str:
        """``var(--x, fallback)`` → 실제 값. 무한루프 방지 depth cap 4."""
        if _depth >= 4:
            return v.strip()
        m = re.match(r"\s*var\s*\(\s*--([\w-]+)\s*(?:,\s*([^)]+))?\s*\)\s*$",
                     v.strip())
        if m:
            ref = m.group(1).lower()
            fallback = (m.group(2) or "").strip()
            if ref in variables:
                return _resolve_value(variables[ref], _depth + 1)
            return fallback
        return v.strip()

    for key, aliases in _CSS_VAR_ALIASES.items():
        for alias in aliases:
            if alias in variables:
                out[key] = _resolve_value(variables[alias])
                break

    # 2. 클래스 / 태그 속성 — 변수에서 못 찾은 키 보강
    def _extract(sel_re: str, prop: str) -> str | None:
        rule_re = re.compile(
            rf"(^|[\s,{{}}]){sel_re}\s*(?:,[^{{]*)?\{{([^}}]+)\}}",
            re.IGNORECASE | re.DOTALL | re.MULTILINE,
        )
        for m in rule_re.finditer(text):
            body = m.group(2)
            pm = re.search(rf"\b{prop}\s*:\s*([^;]+);", body, re.IGNORECASE)
            if pm:
                return _resolve_value(pm.group(1))
        return None

    # button bg/text
    if "button_bg" not in out:
        v = (_extract(r"\.btn-primary", "background-color")
             or _extract(r"\.btn-primary", "background")
             or _extract(r"\.button-primary", "background-color")
             or _extract(r"\.btn\.primary", "background-color"))
        if v:
            out["button_bg"] = v
    if "button_text" not in out:
        v = (_extract(r"\.btn-primary", "color")
             or _extract(r"\.button-primary", "color")
             or _extract(r"\.btn\.primary", "color"))
        if v:
            out["button_text"] = v
    # button shape
    if "button_shape" not in out:
        radius = (_extract(r"\.btn", "border-radius")
                  or _extract(r"\.button", "border-radius"))
        if radius:
            rm = re.search(r"(\d+(?:\.\d+)?)", radius)
            if rm:
                out["button_shape"] = "rounded" if float(rm.group(1)) >= 2 else "square"
    # table header
    if "table_header_bg" not in out:
        v = (_extract(r"thead", "background-color")
             or _extract(r"thead", "background")
             or _extract(r"\.table\s+thead", "background-color")
             or _extract(r"\.grid-header", "background-color")
             or _extract(r"\.ag-header", "background-color"))
        if v:
            out["table_header_bg"] = v
    if "table_header_text" not in out:
        v = (_extract(r"thead", "color")
             or _extract(r"\.grid-header", "color")
             or _extract(r"\.ag-header", "color"))
        if v:
            out["table_header_text"] = v
    # panel bg/border
    if "panel_bg" not in out:
        v = (_extract(r"\.panel", "background-color")
             or _extract(r"\.search-panel", "background-color")
             or _extract(r"\.card", "background-color"))
        if v:
            out["panel_bg"] = v
    if "panel_border" not in out:
        v = _extract(r"\.panel", "border-color") or _extract(r"\.card", "border-color")
        if v:
            out["panel_border"] = v
    # input bg/border
    if "input_bg" not in out:
        v = (_extract(r"input", "background-color")
             or _extract(r"\.form-control", "background-color"))
        if v:
            out["input_bg"] = v
    if "input_border" not in out:
        v = (_extract(r"input", "border-color")
             or _extract(r"\.form-control", "border-color"))
        if v:
            out["input_border"] = v
    # font-family — body / :root / html
    if "font_family" not in out:
        v = (_extract(r"body", "font-family")
             or _extract(r":root", "font-family")
             or _extract(r"html", "font-family"))
        if v:
            out["font_family"] = v

    # font-family 후처리 — 첫 폰트만 + quotes / fallback chain 제거.
    if "font_family" in out:
        ff = out["font_family"].split(",")[0].strip().strip("'\"")
        if ff:
            out["font_family"] = ff

    return out


# ── 소스 매칭 ────────────────────────────────────────────────────────

_SOURCE_EXTS = (".tsx", ".jsx", ".ts", ".js", ".vue")
_SOURCE_EXCLUDE_SEGMENTS = (
    "node_modules", "/dist/", "\\dist\\",
    "/build/", "\\build\\", "/.next/", "\\.next\\",
    "/coverage/", "/.git/", "/__tests__/",
    ".test.", ".spec.",               # 테스트 파일 매칭 회피
    ".stories.", ".story.",           # storybook
    ".d.ts",                          # 타입 선언 파일
)
_SOURCE_SCREEN_DIR_HINTS = ("/pages/", "/screens/", "/views/", "/routes/")
_SOURCE_MAX_CHARS = 8000          # VLM 컨텍스트 보호용 상한
_SOURCE_MIN_SCORE = 4             # 이 이상이어야 매칭으로 인정 (오매칭 방지)
_SOURCE_TOP_N_DIAG = 3            # 매칭 실패 시 진단용으로 보존할 상위 후보 수
_TOKEN_RE = re.compile(r"[^a-zA-Z0-9가-힣]+")


def _normalize_tokens(s: str) -> list[str]:
    """alphanumeric (+ 한글) 토큰을 소문자로 분리."""
    return [t for t in _TOKEN_RE.split(s.lower()) if t]


def _build_source_index(frontend_dir: Path | None) -> list[Path]:
    """frontend_dir 하위의 React/Vue/TS 파일을 1회 수집. 없으면 빈 리스트."""
    if frontend_dir is None or not frontend_dir.is_dir():
        return []
    files: list[Path] = []
    for ext in _SOURCE_EXTS:
        files.extend(frontend_dir.rglob(f"*{ext}"))
    out = []
    for p in files:
        sp = str(p).replace("\\", "/").lower()
        if any(seg.lower() in sp for seg in _SOURCE_EXCLUDE_SEGMENTS):
            continue
        out.append(p)
    return out


def _score_match(stem_tokens: list[str], cand: Path) -> int:
    """캡처 stem 토큰과 소스 파일 path/basename 의 매치 점수.
    - 전체 stem 이 basename 의 substring 이면 +10
    - 각 토큰이 basename 에 있으면 +3, path 에만 있으면 +1
    - 파일이 pages/screens/views/routes 하위면 +2 (화면 컴포넌트 우대)
    """
    if not stem_tokens:
        return 0
    base = cand.stem.lower()
    path_text = str(cand.parent).replace("\\", "/").lower() + "/" + base
    full = "".join(stem_tokens)
    score = 0
    if full and full in re.sub(_TOKEN_RE, "", base):
        score += 10
    for t in stem_tokens:
        if len(t) < 2:
            continue
        if t in base:
            score += 3
        elif t in path_text:
            score += 1
    if any(hint in path_text for hint in _SOURCE_SCREEN_DIR_HINTS):
        score += 2
    return score


def _rank_sources(asis_stem: str,
                  source_index: list[Path]) -> list[tuple[Path, int]]:
    """캡처 stem 에 대한 소스 후보 (path, score) 정렬 리스트 (점수 내림차순)."""
    stem_tokens = _normalize_tokens(asis_stem)
    if not stem_tokens or not source_index:
        return []
    scored = [(cand, _score_match(stem_tokens, cand)) for cand in source_index]
    scored = [(p, s) for p, s in scored if s > 0]
    scored.sort(key=lambda ps: -ps[1])
    return scored


def _match_source(asis_stem: str, source_index: list[Path],
                  mapping: dict[str, Path] | None = None
                  ) -> tuple[Path | None, list[tuple[Path, int]]]:
    """매핑 우선 → 휴리스틱 fallback. (best_path|None, top_candidates) 반환.

    - 명시 매핑(mapping[asis_stem]) 이 있으면 그걸 사용 (점수 무시).
    - 없으면 _rank_sources 의 top1 점수가 임계 이상이면 채택.
    - top_candidates 는 진단용 — 매칭 실패해도 상위 N개 보존.
    """
    if mapping and asis_stem in mapping:
        return mapping[asis_stem], []
    ranked = _rank_sources(asis_stem, source_index)
    top_n = ranked[:_SOURCE_TOP_N_DIAG]
    if not ranked:
        return None, top_n
    best_path, best_score = ranked[0]
    if best_score >= _SOURCE_MIN_SCORE:
        return best_path, top_n
    return None, top_n


def _load_source_mapping(mapping_path: Path | None,
                         frontend_dir: Path | None) -> dict[str, Path]:
    """수기 매핑 YAML 로드. {capture_stem: Path}.

    YAML 값이 절대 경로면 그대로, 상대 경로면 frontend_dir 기준으로
    resolve. 존재하지 않는 파일은 경고 후 스킵.
    """
    if mapping_path is None:
        return {}
    if not mapping_path.is_file():
        raise SystemExit(f"--source-mapping 파일 없음: {mapping_path}")
    try:
        import yaml
    except ImportError as e:
        raise SystemExit("PyYAML 필요 — `pip install pyyaml`") from e
    raw = yaml.safe_load(mapping_path.read_text(encoding="utf-8")) or {}
    if not isinstance(raw, dict):
        raise SystemExit(
            f"--source-mapping YAML 은 {{capture_stem: 경로}} dict 여야 함: "
            f"{mapping_path}"
        )
    out: dict[str, Path] = {}
    for stem, rel in raw.items():
        if rel is None:
            continue
        p = Path(str(rel))
        if not p.is_absolute() and frontend_dir is not None:
            p = frontend_dir / p
        if not p.is_file():
            logger.warning("매핑 파일 미존재 — 스킵: %s → %s", stem, p)
            continue
        out[str(stem)] = p
    return out


def _read_source_snippet(path: Path, max_chars: int = _SOURCE_MAX_CHARS) -> str:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logger.warning("소스 읽기 실패 (%s): %s", path, e)
        return ""
    if len(text) > max_chars:
        text = text[:max_chars] + "\n/* ... (이하 생략) */"
    return text


# 화면 단위 closure (entry + import 그래프 BFS) — analyze-legacy 의
# legacy_react_closure 인프라 재사용. tree-sitter 의존, 미설치 시 미사용.
_CLOSURE_TOKEN_BUDGET = 20000          # 8000 char single-file 보다 훨씬 큼
_CLOSURE_MAX_DEPTH = 3                 # build_closure 기본값

# 한 번만 import 시도, 결과 (모듈 또는 None) 를 캐시
_CLOSURE_IMPORT_TRIED = False
_CLOSURE_MOD: Any = None
_CLOSURE_IMPORT_ERR: str | None = None


def _load_closure_module():
    """legacy_react_closure 를 lazy import. tree-sitter 없으면 None 캐시."""
    global _CLOSURE_IMPORT_TRIED, _CLOSURE_MOD, _CLOSURE_IMPORT_ERR
    if _CLOSURE_IMPORT_TRIED:
        return _CLOSURE_MOD
    _CLOSURE_IMPORT_TRIED = True
    try:
        from . import legacy_react_closure as mod
        _CLOSURE_MOD = mod
    except Exception as e:  # noqa: BLE001 — tree-sitter 미설치 등
        _CLOSURE_IMPORT_ERR = f"{type(e).__name__}: {e}"
        _CLOSURE_MOD = None
    return _CLOSURE_MOD


def _bundle_source_closure(entry_file: Path,
                           frontend_dir: Path) -> tuple[str, dict] | None:
    """entry_file 부터 import 그래프 BFS → Markdown 직렬화 + 메타.

    내부적으로 `legacy_react_closure.build_closure` 호출. tree-sitter
    의존, 미설치/실패 시 None. 반환: (markdown_text, stats_dict).
    """
    mod = _load_closure_module()
    if mod is None:
        return None
    try:
        closure = mod.build_closure(
            entry_file=str(entry_file),
            repo_root=str(frontend_dir),
            patterns=None,
            max_depth=_CLOSURE_MAX_DEPTH,
            token_budget=_CLOSURE_TOKEN_BUDGET,
            verbose=False,
        )
        markdown = mod.serialize_for_llm(closure)
    except Exception as e:  # noqa: BLE001 — closure 실패 시 fallback
        logger.warning("closure 빌드 실패 (%s): %s — 단일 파일 fallback",
                       entry_file, e)
        return None
    stats = {
        "entry_name": closure.entry_name,
        "file_count": len(closure.files),
        "total_tokens": closure.total_tokens,
        "truncated": closure.truncated,
        "api_calls": len(closure.api_calls),
        "popup_refs": len(closure.popup_refs),
        "skipped_external": len(closure.skipped_external),
        "files": [
            {"rel_path": f.rel_path, "depth": f.depth, "mode": f.mode}
            for f in closure.files
        ],
    }
    return markdown, stats


# ── LLM 추출 ────────────────────────────────────────────────────────


def extract_layout(asis_image: Path, template_images: list[Path],
                   config: dict,
                   source_path: Path | None = None,
                   frontend_dir: Path | None = None
                   ) -> tuple[dict, dict | None]:
    """VLM 1회 호출 → (layout dict, source_attachment_meta | None).

    source_path 가 주어지면 React 소스를 프롬프트 하단에 첨부:
    - frontend_dir 동시 주어지고 tree-sitter 설치돼있으면 `build_closure`
      로 entry+import 그래프 BFS 번들 (여러 파일, ~20K tokens).
    - 아니면 단일 파일 8000 char fallback.

    반환의 2번째 요소는 진단용 메타 (mode='closure'|'single'|None, file_count
    등). LLM 응답이 dict 아니면 layout 은 빈 dict.
    """
    prompt = _USER_PROMPT
    attach_meta: dict | None = None
    if source_path is not None:
        bundle = (_bundle_source_closure(source_path, frontend_dir)
                  if frontend_dir is not None else None)
        if bundle is not None:
            md, stats = bundle
            prompt = (
                _USER_PROMPT
                + "\n\n=== React 소스 첨부 (closure bundle, "
                f"{stats['file_count']} 파일, ~{stats['total_tokens']} tokens) ===\n"
                + "이 화면의 entry 컴포넌트 + import 로 따라간 자식 컴포넌트"
                + " 들이 모두 포함되어 있습니다.\n"
                + "이 소스가 search_fields / table_columns / buttons 의 정답.\n"
                + "이미지는 regions (위치/크기) 와 page_title 추론에만 사용.\n\n"
                + md
            )
            print(f"    소스 첨부 (closure): {stats['file_count']} 파일, "
                  f"~{stats['total_tokens']} tokens, "
                  f"entry={stats['entry_name']}"
                  f"{', truncated' if stats['truncated'] else ''}")
            attach_meta = {"mode": "closure", **stats}
        else:
            snippet = _read_source_snippet(source_path)
            if snippet:
                rel = source_path.name
                prompt = (
                    _USER_PROMPT
                    + f"\n\n=== React 소스 첨부 ({rel}, single file) ===\n"
                    + "이 소스가 search_fields / table_columns / buttons 의 정답.\n"
                    + "이미지는 regions (위치/크기) 와 page_title 추론에만 사용.\n\n"
                    + snippet
                )
                hint = (f" (closure 미사용: {_CLOSURE_IMPORT_ERR})"
                        if _CLOSURE_IMPORT_ERR else "")
                print(f"    소스 첨부 (single): {rel} ({len(snippet)} chars){hint}")
                attach_meta = {"mode": "single", "chars": len(snippet),
                               "file": str(source_path)}

    image_paths = [str(asis_image)] + [str(t) for t in template_images]
    result = _call_llm(
        prompt=prompt,
        config=config,
        label=f"screen_{asis_image.stem}",
        system_prompt=_SYSTEM_PROMPT,
        image_paths=image_paths,
    )
    if not isinstance(result, dict):
        logger.warning("LLM 응답이 dict 가 아님: %s — 빈 레이아웃으로 대체", asis_image.name)
        return {}, attach_meta
    return result, attach_meta


# ── VLM 으로 HTML 추출 (옵션 D: --export-html) ───────────────────────

_HTML_SYSTEM_PROMPT = (
    "당신은 차세대 SI 프론트엔드 개발자입니다. AS-IS 화면 캡처와 (있으면) "
    "TO-BE 디자인 템플릿 캡처를 보고, 첨부된 TO-BE CSS 의 클래스/규칙을 "
    "정확히 사용하는 HTML 마크업 (body 안 내용만) 을 생성하세요. "
    "결과는 HTML 코드 블록 하나만, 설명/마크다운 텍스트/주석 금지."
)

_HTML_PROMPT_TEMPLATE = """다음 입력으로 TO-BE HTML body 안 마크업을 생성:

[AS-IS] 첫 번째 이미지 = 변환 대상 화면 캡처
[Template] 그 뒤 이미지들 = TO-BE 디자인 시각 참조
[CSS] 첨부된 TO-BE CSS = 사용 가능한 클래스/selector 의 출처
[Source] (있으면) AS-IS React 소스 = 라벨/컬럼/버튼 텍스트 정답

규칙:
1. **반드시 CSS 안에 정의된 class/selector 만 사용**. 임의로 새 클래스
   이름 발명 금지. CSS 에 `.btn-primary` `.search-form .field-row`
   같은 게 있으면 그걸 그대로 써라.
2. 구조도 CSS 가 기대하는 nesting 그대로 (`.form > .field > input` 등).
3. AS-IS 의 모든 라벨/컬럼/버튼 텍스트를 빠짐없이 옮긴다. React 소스가
   첨부됐으면 그 텍스트가 정답. 이미지에서만 다르게 보여도 소스 우선.
4. 출력은 `<body>` 안 들어갈 마크업만. `<html>` / `<head>` /
   `<link rel=stylesheet>` 같은 wrap 은 생성하지 마라 (호출 측이 추가).
5. 결과를 코드 블록으로 감싸 출력:
   ```html
   <div class=...>
     ...
   </div>
   ```

=== TO-BE CSS (사용 가능한 selector 목록) ===
{css_text}
"""


def _truncate_css_for_prompt(css_text: str, max_chars: int = 12000) -> str:
    """CSS 가 너무 크면 VLM 컨텍스트 안 넘어가도록 잘라낸다.
    selector + property name 만 추출하는 정밀화는 후속.
    """
    if len(css_text) <= max_chars:
        return css_text
    return css_text[:max_chars] + "\n/* ... (이하 생략 — CSS truncated) */"


_HTML_FENCE_RE = re.compile(r"```(?:html|HTML)?\s*([\s\S]*?)```", re.MULTILINE)


def _strip_html_fence(text: str) -> str:
    """LLM 응답에서 ```html ...``` 펜스 안 본문만 추출. 없으면 원본."""
    if not text:
        return ""
    m = _HTML_FENCE_RE.search(text)
    if m:
        return m.group(1).strip()
    # 펜스 없으면 trim 후 그대로 (단, 첫 < 부터 마지막 > 까지)
    text = text.strip()
    first = text.find("<")
    last = text.rfind(">")
    if first != -1 and last != -1 and last > first:
        return text[first:last + 1]
    return text


def extract_html(asis_image: Path, template_images: list[Path],
                 css_text: str, source_text: str | None,
                 config: dict) -> str:
    """VLM 1회 호출 → TO-BE HTML body 마크업 문자열. 실패 시 빈 문자열.

    `_call_llm` 가 JSON 응답을 기대해서 여기서는 직접 OpenAI client 호출
    경로를 쓰지 않고, 텍스트 응답을 받기 위해 chat completions API 를
    직접 사용한다 (CSS + 이미지 묶음).
    """
    # 직접 OpenAI 호환 chat call (JSON 강제 안 함)
    try:
        from openai import OpenAI
    except ImportError as e:
        raise SystemExit("openai 패키지 필요 — `pip install openai`") from e

    llm_config = config.get("llm", {})
    api_key = (os.environ.get("PATTERN_LLM_API_KEY")
               or os.environ.get("LLM_API_KEY")
               or llm_config.get("api_key", "ollama"))
    api_base = (os.environ.get("PATTERN_LLM_API_BASE")
                or os.environ.get("LLM_API_BASE")
                or llm_config.get("api_base", "http://localhost:11434/v1"))
    model = (os.environ.get("PATTERN_LLM_MODEL")
             or os.environ.get("LLM_MODEL")
             or llm_config.get("model", "llama3"))

    client = OpenAI(api_key=api_key, base_url=api_base)

    # multimodal user content (AS-IS + templates)
    import base64
    user_content: list[dict] = []
    user_content.append({
        "type": "text",
        "text": _HTML_PROMPT_TEMPLATE.format(
            css_text=_truncate_css_for_prompt(css_text)
        ),
    })
    for img in [asis_image] + list(template_images):
        try:
            data = Path(img).read_bytes()
            b64 = base64.b64encode(data).decode("ascii")
            ext = Path(img).suffix.lower().lstrip(".") or "png"
            mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
            user_content.append({
                "type": "image_url",
                "image_url": {"url": f"data:{mime};base64,{b64}"},
            })
        except Exception as e:
            logger.warning("HTML extract — 이미지 base64 실패 (%s): %s", img, e)
    if source_text:
        user_content.append({
            "type": "text",
            "text": (
                "\n=== AS-IS React 소스 (라벨/컬럼/버튼 텍스트 정답) ===\n"
                + source_text[:_SOURCE_MAX_CHARS]
            ),
        })

    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": _HTML_SYSTEM_PROMPT},
                {"role": "user", "content": user_content},
            ],
            temperature=0.1,
            timeout=300,
        )
        raw = response.choices[0].message.content or ""
    except Exception as e:
        logger.error("HTML extract LLM 호출 실패 (%s): %s", asis_image.name, e)
        return ""
    return _strip_html_fence(raw)


def render_html(screens: list[tuple[str, str]],
                css_text: str,
                output_dir: Path,
                css_filename: str = "tobe_style.css") -> None:
    """screens=[(stem, html_body), ...] → 화면별 html + index.html.

    - 같은 폴더에 css 파일 복사 (`<link rel=stylesheet href=...>`)
    - index.html: 모든 화면 링크 목록
    - 각 화면 html: 단순 wrap (<!DOCTYPE> + <head> 링크 + body)
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    css_path = output_dir / css_filename
    css_path.write_text(css_text, encoding="utf-8")

    rendered = 0
    for stem, body in screens:
        if not body:
            continue
        page = (
            "<!DOCTYPE html>\n"
            "<html lang=\"ko\">\n"
            "<head>\n"
            "  <meta charset=\"utf-8\">\n"
            f"  <title>{_html_escape(stem)}</title>\n"
            f"  <link rel=\"stylesheet\" href=\"{css_filename}\">\n"
            "</head>\n"
            "<body>\n"
            f"{body}\n"
            "</body>\n"
            "</html>\n"
        )
        (output_dir / f"{stem}.html").write_text(page, encoding="utf-8")
        rendered += 1

    # index
    items = "\n".join(
        f"  <li><a href=\"{_html_escape(stem)}.html\">{_html_escape(stem)}</a></li>"
        for stem, body in screens if body
    )
    index = (
        "<!DOCTYPE html>\n"
        "<html lang=\"ko\">\n"
        "<head>\n"
        "  <meta charset=\"utf-8\">\n"
        "  <title>TO-BE 화면 인덱스</title>\n"
        f"  <link rel=\"stylesheet\" href=\"{css_filename}\">\n"
        "  <style>body{font-family:sans-serif;padding:24px}"
        "ul{line-height:1.8}</style>\n"
        "</head>\n"
        "<body>\n"
        f"  <h1>TO-BE 화면 인덱스 ({rendered}장)</h1>\n"
        "  <ul>\n"
        f"{items}\n"
        "  </ul>\n"
        "</body>\n"
        "</html>\n"
    )
    (output_dir / "index.html").write_text(index, encoding="utf-8")


def _html_escape(s: str) -> str:
    return (str(s).replace("&", "&amp;").replace("<", "&lt;")
            .replace(">", "&gt;").replace("\"", "&quot;"))


# ── PPTX 렌더 ────────────────────────────────────────────────────────


def _slide_dims_cm(prs) -> tuple[float, float]:
    """슬라이드 가로/세로 cm. EMU 가 (1cm = 360000 EMU) 정확."""
    from pptx.util import Cm, Emu
    one_cm = int(Cm(1))
    return (prs.slide_width / one_cm, prs.slide_height / one_cm)


def _clamp01(v) -> float | None:
    try:
        f = float(v)
    except (TypeError, ValueError):
        return None
    if f != f:  # NaN
        return None
    return max(0.0, min(1.0, f))


def _resolve_bbox(region: dict | None, default_cm: tuple[float, float, float, float],
                  prs) -> tuple[float, float, float, float]:
    """region 의 normalized bbox → cm 튜플. 부재/오류 시 default_cm 그대로."""
    if not isinstance(region, dict):
        return default_cm
    x = _clamp01(region.get("x"))
    y = _clamp01(region.get("y"))
    w = _clamp01(region.get("w"))
    h = _clamp01(region.get("h"))
    if None in (x, y, w, h) or w <= 0 or h <= 0:
        return default_cm
    # clamp w/h so x+w, y+h stay within slide
    w = min(w, 1.0 - x)
    h = min(h, 1.0 - y)
    sw, sh = _slide_dims_cm(prs)
    return (x * sw, y * sh, w * sw, h * sh)


def _detect_slide_aspect_inches(template_image: Path | None) -> tuple[float, float]:
    """첫 템플릿 이미지의 가로:세로 비율 → 슬라이드 가로/세로 (inch).
    >= 1.5 → 16:9 (13.333 x 7.5), 미만 → 4:3 (10.0 x 7.5).
    이미지 없거나 읽기 실패면 16:9 기본.
    """
    if template_image is None or not Path(template_image).is_file():
        return (13.333, 7.5)
    try:
        from PIL import Image
        with Image.open(template_image) as img:
            w, h = img.size
        if h <= 0:
            return (13.333, 7.5)
        return (13.333, 7.5) if (w / h) >= 1.5 else (10.0, 7.5)
    except Exception:
        return (13.333, 7.5)


def _draw_title(slide, text: str, prs, region: dict | None = None,
                style: dict | None = None) -> None:
    from pptx.util import Cm, Pt
    style = _resolve_style(style)
    sw, _ = _slide_dims_cm(prs)
    default = (0.8, 0.6, sw - 1.6, 1.2)
    x, y, w, h = _resolve_bbox(region, default, prs)
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or "(제목 없음)"
    r = p.runs[0]
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.name = style["font_family"]
    r.font.color.rgb = _parse_hex_color(style.get("title_color"), (0x1F, 0x3A, 0x5F))


def _draw_section_label(slide, text: str, x_cm: float, y_cm: float,
                        style: dict, w_cm: float = 6.0) -> None:
    from pptx.util import Cm, Pt
    tb = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.name = style["font_family"]
    r.font.color.rgb = _parse_hex_color(style.get("section_label_color"),
                                        (0x44, 0x44, 0x44))


def _draw_search(slide, fields: list[dict], prs,
                 region: dict | None = None,
                 fallback_top_cm: float = 2.2,
                 style: dict | None = None) -> float:
    """검색 패널. region 있으면 bbox 사용, 없으면 fallback_top_cm 기준 스택.
    반환: 다음 섹션이 사용할 누적 높이 (cm). region 사용 시 0.0."""
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE

    if not fields:
        return 0.0

    style = _resolve_style(style)
    sw, _ = _slide_dims_cm(prs)
    panel_bg = _parse_hex_color(style.get("panel_bg"), (0xF5, 0xF7, 0xFA))
    panel_border = _parse_hex_color(style.get("panel_border"), (0xCF, 0xD6, 0xE0))
    input_bg = _parse_hex_color(style.get("input_bg"), (0xFF, 0xFF, 0xFF))
    input_border = _parse_hex_color(style.get("input_border"), (0xA0, 0xA8, 0xB4))
    field_label_color = _parse_hex_color(style.get("field_label_color"),
                                         (0x33, 0x33, 0x33))
    input_placeholder_color = _parse_hex_color(
        style.get("input_placeholder_color"), (0x99, 0x99, 0x99))
    font_name = style["font_family"]

    if region:
        cols = max(1, int(region.get("cols") or 3))
        default = (0.8, fallback_top_cm + 0.7, sw - 1.6, 2.0)
        panel_x, panel_y, panel_w, panel_h = _resolve_bbox(region, default, prs)
        label_y = max(0.0, panel_y - 0.7)
        _draw_section_label(slide, "검색 조건", panel_x, label_y, style)
        advance = 0.0
    else:
        cols = 3
        rows = (len(fields) + cols - 1) // cols
        panel_x = 0.8
        panel_w = sw - 1.6
        cell_h = 1.0
        panel_h = rows * cell_h + 0.3
        _draw_section_label(slide, "검색 조건", panel_x, fallback_top_cm, style)
        panel_y = fallback_top_cm + 0.7
        advance = panel_h + 0.7

    rows = (len(fields) + cols - 1) // cols
    cell_w = panel_w / cols
    cell_h = panel_h / rows if rows > 0 else panel_h

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(panel_x), Cm(panel_y),
                                Cm(panel_w), Cm(panel_h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = panel_bg
    bg.line.color.rgb = panel_border
    bg.line.width = Pt(0.5)
    if bg.has_text_frame:
        bg.text_frame.text = ""

    for idx, f in enumerate(fields):
        r_idx, c_idx = divmod(idx, cols)
        x = panel_x + c_idx * cell_w + 0.2
        y = panel_y + 0.15 + r_idx * cell_h
        usable_w = cell_w - 0.4
        lb = slide.shapes.add_textbox(Cm(x), Cm(y),
                                      Cm(usable_w * 0.4), Cm(0.6))
        lp = lb.text_frame.paragraphs[0]
        lp.text = (f.get("label") or "").strip() or "(라벨)"
        lp.runs[0].font.size = Pt(8)
        lp.runs[0].font.name = font_name
        lp.runs[0].font.color.rgb = field_label_color
        ib = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(x + usable_w * 0.4 + 0.1), Cm(y + 0.05),
            Cm(usable_w * 0.55), Cm(min(0.55, cell_h - 0.3)),
        )
        ib.fill.solid()
        ib.fill.fore_color.rgb = input_bg
        ib.line.color.rgb = input_border
        ib.line.width = Pt(0.5)
        tf = ib.text_frame
        tf.margin_left = Cm(0.1)
        tf.margin_right = Cm(0.1)
        tp = tf.paragraphs[0]
        tp.text = f"<{(f.get('type') or 'text').strip()}>"
        tp.runs[0].font.size = Pt(7)
        tp.runs[0].font.name = font_name
        tp.runs[0].font.color.rgb = input_placeholder_color

    return advance


def _draw_table(slide, columns: list[str], prs,
                region: dict | None = None,
                fallback_top_cm: float = 0.0,
                style: dict | None = None) -> float:
    """표 헤더 + 빈 행 3줄. region 있으면 bbox 사용, 없으면 fallback 스택."""
    from pptx.util import Cm, Pt

    if not columns:
        return 0.0

    style = _resolve_style(style)
    sw, _ = _slide_dims_cm(prs)
    header_bg = _parse_hex_color(style.get("table_header_bg"), (0x1F, 0x3A, 0x5F))
    header_text = _parse_hex_color(style.get("table_header_text"),
                                   (0xFF, 0xFF, 0xFF))
    row_bg = _parse_hex_color(style.get("table_row_bg"), (0xFF, 0xFF, 0xFF))
    font_name = style["font_family"]

    if region:
        default = (0.8, fallback_top_cm + 0.7, sw - 1.6, 5.0)
        tbl_x, tbl_y, tbl_w, tbl_h = _resolve_bbox(region, default, prs)
        label_y = max(0.0, tbl_y - 0.7)
        _draw_section_label(slide, "조회 결과", tbl_x, label_y, style)
        advance = 0.0
    else:
        n_rows_default = 4
        tbl_x = 0.8
        tbl_w = sw - 1.6
        tbl_h = 0.7 * n_rows_default
        _draw_section_label(slide, "조회 결과", tbl_x, fallback_top_cm, style)
        tbl_y = fallback_top_cm + 0.7
        advance = tbl_h + 0.7

    n_rows = 4  # header + 3 empty
    n_cols = max(len(columns), 1)

    shape = slide.shapes.add_table(n_rows, n_cols,
                                   Cm(tbl_x), Cm(tbl_y),
                                   Cm(tbl_w), Cm(tbl_h))
    table = shape.table

    for ci, col in enumerate(columns):
        cell = table.cell(0, ci)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.name = font_name
                r.font.color.rgb = header_text

    for ri in range(1, n_rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg

    return advance


def _draw_buttons(slide, buttons: list[str], prs,
                  region: dict | None = None,
                  fallback_top_cm: float = 0.0,
                  style: dict | None = None) -> float:
    from pptx.util import Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE

    if not buttons:
        return 0.0

    style = _resolve_style(style)
    sw, _ = _slide_dims_cm(prs)
    btn_bg = _parse_hex_color(style.get("button_bg"), (0x1F, 0x3A, 0x5F))
    btn_text = _parse_hex_color(style.get("button_text"), (0xFF, 0xFF, 0xFF))
    font_name = style["font_family"]
    shape_kind = (style.get("button_shape") or "rounded").strip().lower()
    btn_shape_enum = (MSO_SHAPE.RECTANGLE if shape_kind == "square"
                      else MSO_SHAPE.ROUNDED_RECTANGLE)

    if region:
        default = (sw - 8.0, fallback_top_cm, 7.5, 0.8)
        bar_x, bar_y, bar_w, bar_h = _resolve_bbox(region, default, prs)
        align = (region.get("align") or "right").lower()
        advance = 0.0
    else:
        bar_x = 0.8
        bar_y = fallback_top_cm
        bar_w = sw - 1.6
        bar_h = 0.8
        align = "right"
        advance = bar_h + 0.3

    n = len(buttons)
    gap = 0.25
    btn_w = min(2.4, (bar_w - (n - 1) * gap) / max(n, 1))
    btn_w = max(btn_w, 1.5)
    btn_h = min(bar_h, 0.9)
    total_w = n * btn_w + (n - 1) * gap

    if align == "left":
        x0 = bar_x
    elif align == "center":
        x0 = bar_x + (bar_w - total_w) / 2
    else:  # right
        x0 = bar_x + bar_w - total_w

    y0 = bar_y + (bar_h - btn_h) / 2

    for i, label in enumerate(buttons):
        x = x0 + i * (btn_w + gap)
        shape = slide.shapes.add_shape(
            btn_shape_enum,
            Cm(x), Cm(y0), Cm(btn_w), Cm(btn_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = btn_bg
        shape.line.color.rgb = btn_bg
        tf = shape.text_frame
        p = tf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
        p.text = str(label)
        r = p.runs[0]
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.name = font_name
        r.font.color.rgb = btn_text

    return advance


def _draw_notes(slide, text: str, prs,
                region: dict | None = None,
                fallback_top_cm: float = 0.0,
                style: dict | None = None) -> None:
    if not text:
        return
    from pptx.util import Cm, Pt
    style = _resolve_style(style)
    sw, _ = _slide_dims_cm(prs)
    default = (0.8, fallback_top_cm, sw - 1.6, 1.5)
    x, y, w, h = _resolve_bbox(region, default, prs)
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"※ {text}"
    r = p.runs[0]
    r.font.size = Pt(7)
    r.font.italic = True
    r.font.name = style["font_family"]
    r.font.color.rgb = _parse_hex_color(style.get("notes_color"),
                                        (0x55, 0x55, 0x55))


def render_pptx(layouts: list[tuple[str, dict]], output_path: Path,
                aspect_hint_image: Path | None = None,
                style: dict | None = None) -> None:
    """layout 리스트 → PPTX. 슬라이드당 1화면.

    aspect_hint_image 가 주어지면 그 이미지의 가로:세로 비율로 슬라이드
    크기를 설정 (16:9 or 4:3). 미지정 시 16:9 기본.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        raise SystemExit(
            "python-pptx 미설치. `pip install python-pptx` 또는 폐쇄망 wheel "
            "설치 (requirements.txt 주석 참고)"
        ) from e

    prs = Presentation()
    w_in, h_in = _detect_slide_aspect_inches(aspect_hint_image)
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    blank_layout = prs.slide_layouts[6]

    style = _resolve_style(style)

    for screen_name, layout in layouts:
        slide = prs.slides.add_slide(blank_layout)
        regions = layout.get("regions") or {}

        title = (layout.get("page_title") or "").strip() or screen_name
        _draw_title(slide, title, prs, region=regions.get("title"),
                    style=style)

        # regions 가 있으면 모두 절대 좌표, 없으면 cursor 스택
        if regions:
            _draw_search(slide, layout.get("search_fields") or [], prs,
                         region=regions.get("search_panel"), style=style)
            _draw_table(slide, layout.get("table_columns") or [], prs,
                        region=regions.get("table"), style=style)
            _draw_buttons(slide, layout.get("buttons") or [], prs,
                          region=regions.get("buttons"), style=style)
            _draw_notes(slide, (layout.get("notes") or "").strip(), prs,
                        region=regions.get("notes"), style=style)
        else:
            cursor = 2.2
            cursor += _draw_search(slide, layout.get("search_fields") or [],
                                   prs, fallback_top_cm=cursor, style=style)
            cursor += _draw_table(slide, layout.get("table_columns") or [],
                                  prs, fallback_top_cm=cursor, style=style)
            cursor += _draw_buttons(slide, layout.get("buttons") or [],
                                    prs, fallback_top_cm=cursor, style=style)
            _draw_notes(slide, (layout.get("notes") or "").strip(),
                        prs, fallback_top_cm=cursor, style=style)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
    except PermissionError as e:
        raise SystemExit(
            f"PPTX 저장 실패 — {output_path}\n"
            f"  원인: 파일이 다른 프로그램(PowerPoint 등)에 열려있어 잠긴 상태.\n"
            f"  해결: 1) 해당 PowerPoint 창 닫고 재실행, 또는\n"
            f"        2) --output <다른경로.pptx> 로 다른 파일명 지정.\n"
            f"  상세: {e}"
        ) from e


# ── 파이프라인 엔트리 ───────────────────────────────────────────────


def _dump_layout(layout: dict, asis_image: Path, template_images: list[Path],
                 dump_dir: Path,
                 matched_source: Path | None = None,
                 source_candidates: list[tuple[Path, int]] | None = None,
                 source_attachment: dict | None = None
                 ) -> Path:
    """Persist the parsed VLM layout dict for post-hoc inspection.

    렌더링 결과가 기대와 다를 때 모델이 실제로 무엇을 추출했는지 확인
    하기 위한 디버그 산출물. ``_call_llm`` 의 raw text dump 는 JSON 파싱
    실패 때만 떨어지므로, 성공 케이스의 추출 결과는 여기서 따로 저장.
    """
    dump_dir.mkdir(parents=True, exist_ok=True)
    record = {
        "asis_image": asis_image.name,
        "template_images": [t.name for t in template_images],
        "matched_source": str(matched_source) if matched_source else None,
        "source_match_candidates": [
            {"path": str(p), "score": s}
            for p, s in (source_candidates or [])
        ],
        "source_attachment": source_attachment,
        "layout": layout,
    }
    path = dump_dir / f"{asis_image.stem}.json"
    path.write_text(
        json.dumps(record, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return path


def convert(captures_dir: Path, templates_dir: Path | None,
            output_pptx: Path, config: dict,
            frontend_dir: Path | None = None,
            source_mapping_path: Path | None = None,
            style_css_path: Path | None = None,
            export_html: bool = False) -> dict[str, Any]:
    asis_imgs = _list_images(captures_dir)
    tmpl_imgs = _list_images(templates_dir) if templates_dir else []
    if not asis_imgs:
        raise SystemExit(f"AS-IS 캡처 없음 (png/jpg): {captures_dir}")
    if templates_dir and not tmpl_imgs:
        raise SystemExit(f"템플릿 캡처 없음 (png/jpg): {templates_dir}")

    # ``--style-css`` 명시 안 했어도 ``input/tobe_style.css`` 있으면 자동 사용.
    # CWD 또는 repo root (이 모듈 위치 기준) 둘 다 시도.
    if style_css_path is None:
        candidates = [
            Path("input") / "tobe_style.css",  # CWD 기준
            Path(__file__).resolve().parent.parent / "input" / "tobe_style.css",  # repo root 기준
        ]
        for c in candidates:
            if c.is_file():
                style_css_path = c
                break
    elif not Path(style_css_path).is_file():
        # 명시 경로가 잘못된 경우 silent skip 대신 명확히 알림.
        print(f"  ⚠ --style-css 파일 없음: {style_css_path}")
        style_css_path = None

    print(f"  AS-IS 캡처 {len(asis_imgs)}장"
          + (f" / 템플릿 {len(tmpl_imgs)}장 참조" if tmpl_imgs else "")
          + (f" / TO-BE CSS {style_css_path}" if style_css_path else ""))
    if not tmpl_imgs and not style_css_path:
        print("  ⚠ 템플릿 / CSS 둘 다 없음 — default 스타일로 진행")

    source_index = _build_source_index(frontend_dir)
    if frontend_dir is not None:
        if source_index:
            print(f"  프론트 소스 인덱싱: {frontend_dir} ({len(source_index)} 파일)")
        else:
            print(f"  ⚠ 프론트 소스 인덱스 비어있음: {frontend_dir}")

    mapping = _load_source_mapping(source_mapping_path, frontend_dir)
    if mapping:
        print(f"  수기 매핑 로드: {source_mapping_path} ({len(mapping)} 항목)")

    # 스타일 — CSS 우선 (LLM 0), 템플릿이 있으면 VLM 추출로 보강 (CSS 가 우선).
    style_from_css = parse_css_style(style_css_path) if style_css_path else {}
    discovered_vars: list[str] = []
    if style_css_path:
        try:
            _t = Path(style_css_path).read_text(encoding="utf-8", errors="ignore")
            _t = re.sub(r"/\*.*?\*/", "", _t, flags=re.DOTALL)
            discovered_vars = sorted(set(re.findall(r"--([\w-]+)\s*:", _t)))
        except Exception:
            pass
        if style_from_css:
            print(f"  CSS 스타일 파싱 ✓ {len(style_from_css)} 키 적용: "
                  f"{', '.join(sorted(style_from_css))}")
        else:
            print(f"  ⚠ CSS 파싱 0 키 추출 — 변수명/클래스명이 alias 리스트와 "
                  f"불일치 가능성.")
            if discovered_vars:
                preview = ", ".join(f"--{v}" for v in discovered_vars[:12])
                more = (f" (외 {len(discovered_vars) - 12} 개)"
                        if len(discovered_vars) > 12 else "")
                print(f"    발견된 CSS 변수: {preview}{more}")
                print(f"    → 매핑 필요 시 oracle_embeddings/screen_converter.py "
                      f"의 _CSS_VAR_ALIASES 에 변수명 추가")
    style_from_vlm = extract_style_profile(tmpl_imgs, config) if tmpl_imgs else {}
    style_raw = {**style_from_vlm, **style_from_css}   # CSS 가 VLM 덮어쓰기
    style = _resolve_style(style_raw)
    style_path = output_pptx.parent / "style_profile.json"
    style_path.parent.mkdir(parents=True, exist_ok=True)
    style_path.write_text(
        json.dumps({
            "extracted_from_css": style_from_css,
            "discovered_css_variables": discovered_vars,
            "extracted_from_vlm": style_from_vlm,
            "resolved": style,
        }, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"    스타일 저장: {style_path.name} "
          f"({len(style_raw)} 키 추출, {len(style)} 키 적용)")

    dump_dir = output_pptx.parent / "llm_raw"

    layouts: list[tuple[str, dict]] = []
    fail = 0
    matched = 0
    matched_via_mapping = 0
    closure_used = 0
    for img in asis_imgs:
        print(f"  → {img.name}")
        source_path, candidates = _match_source(img.stem, source_index,
                                                mapping=mapping)
        if source_path is not None:
            matched += 1
            if mapping and img.stem in mapping:
                matched_via_mapping += 1
        elif source_index:
            # 휴리스틱 매칭 실패 시 진단: 최고 점수 표시
            top_score = candidates[0][1] if candidates else 0
            print(f"    소스 매칭 실패 (최고 점수 {top_score}/{_SOURCE_MIN_SCORE}, "
                  f"상위 후보는 llm_raw/{img.stem}.json 의 "
                  f"source_match_candidates 참고)")
        layout, attach_meta = extract_layout(
            img, tmpl_imgs, config,
            source_path=source_path,
            frontend_dir=frontend_dir,
        )
        if attach_meta and attach_meta.get("mode") == "closure":
            closure_used += 1
        if not layout:
            fail += 1
        layouts.append((img.stem, layout))
        _dump_layout(layout, img, tmpl_imgs, dump_dir,
                     matched_source=source_path,
                     source_candidates=candidates,
                     source_attachment=attach_meta)

    render_pptx(layouts, output_pptx,
                aspect_hint_image=tmpl_imgs[0] if tmpl_imgs else None,
                style=style)

    # --export-html: 화면별 HTML 생성 (CSS 필수)
    html_stats: dict[str, Any] = {}
    if export_html:
        if not style_css_path or not Path(style_css_path).is_file():
            print("  ⚠ --export-html 지정됐으나 TO-BE CSS 가 없음 — HTML 스킵")
        else:
            html_stats = _generate_html_per_screen(
                asis_imgs, tmpl_imgs, style_css_path, source_index,
                mapping, frontend_dir, config, output_pptx.parent / "html",
            )

    return {
        "total": len(asis_imgs),
        "templates": len(tmpl_imgs),
        "fail": fail,
        "source_matched": matched,
        "source_matched_via_mapping": matched_via_mapping,
        "source_indexed": len(source_index),
        "closure_used": closure_used,
        "style_keys_extracted": len(style_raw),
        "pptx": str(output_pptx),
        "llm_raw_dir": str(dump_dir),
        "style_profile": str(style_path),
        "html_dir": html_stats.get("dir") if html_stats else None,
        "html_generated": html_stats.get("generated") if html_stats else 0,
    }


def _generate_html_per_screen(asis_imgs: list[Path],
                              tmpl_imgs: list[Path],
                              style_css_path: Path,
                              source_index: list[Path],
                              mapping: dict[str, Path],
                              frontend_dir: Path | None,
                              config: dict,
                              output_dir: Path) -> dict[str, Any]:
    """--export-html 옵션 시 호출. 각 캡처별로 VLM 한테 CSS+이미지+소스
    던져서 HTML body 받고, render_html 로 파일로 떨어뜨림.
    """
    print(f"  HTML 생성 (--export-html): VLM 으로 화면별 HTML body 추출")
    try:
        css_text = Path(style_css_path).read_text(encoding="utf-8",
                                                  errors="ignore")
    except Exception as e:
        print(f"  ⚠ CSS 읽기 실패: {e} — HTML 스킵")
        return {}

    screens: list[tuple[str, str]] = []
    generated = 0
    for img in asis_imgs:
        source_path, _ = _match_source(img.stem, source_index, mapping=mapping)
        source_text = ""
        if source_path is not None and frontend_dir is not None:
            bundle = _bundle_source_closure(source_path, frontend_dir)
            if bundle is not None:
                source_text, _ = bundle
            else:
                source_text = _read_source_snippet(source_path)
        print(f"    → {img.name}")
        body = extract_html(img, tmpl_imgs, css_text,
                            source_text or None, config)
        if body:
            generated += 1
        screens.append((img.stem, body))

    render_html(screens, css_text, output_dir,
                css_filename=Path(style_css_path).name)
    return {"dir": str(output_dir), "generated": generated}

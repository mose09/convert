"""Phase C — Screen Layout 추출 + HTML mockup 생성.

각 React 화면 파일을 LLM 으로 분석해 구조화 JSON (Page Title / Search
Panel / DataTable / Edit Mode / Tabs / Events with backend URLs) 추출 후
정적 HTML mockup 으로 렌더. 폐쇄망에서 외부 의존 0 (인라인 CSS).

옵트인 플래그: ``--extract-screen-layout`` (Phase A/B 의존하지 않음 —
``handlers_by_url`` 만 있으면 됨).

별도 옵션: ``--render-screenshots`` (스텁) — Playwright 등으로 진짜 화면
렌더는 사용자 환경에 React 빌드/실행 가능해야 해서 별도 path.

캐시: ``output/legacy_analysis/.screen_cache/<hash>.json``.
"""

from __future__ import annotations

import hashlib
import html
import json
import logging
import os
import re
import shutil
import subprocess
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


SCREEN_SCHEMA_VERSION = "v22"  # v22: <Modal><X/></Modal> 안 컴포넌트 (X) 의 파일은 popup 컨텐츠로 간주, main search panel 에서 제외 — 캐시 무효화

_DEFAULT_CONFIG = {
    "llm_max_chars": 32000,    # 큰 React 파일 대응 (Qwen 397B 컨텍스트 활용)
    "llm_batch_size": 1,       # 화면당 1 호출 (배치 작음)
    "max_screens": 200,
}


@dataclass
class ScreenField:
    label: str = ""
    name: str = ""           # 필드(영문) — input element 의 id 우선 / name / field
    component: str = ""
    default: str = ""
    options: str = ""
    events: str = ""         # "onChange / onClick" 등 (공백 구분)
    required: bool = False   # 필수 여부 (label className required 또는 required prop)
    # 화면정의서 9컬럼 표 양식 (grid 와 parallel)
    placeholder: str = ""    # placeholder attr — UI 가시값 (default 보다 우선)
    max_length: str = ""     # maxLength — keyboard input 자릿수 제한
    input_data_type: str = ""   # String / Number / Date / "" (비입력)
    ui_type: str = ""        # Select(Single) / Text Field(Search Box) / DatePicker 등
    action: str = ""         # 단순 dropdown 은 옵션 줄바꿈 / LLM cascading 판단
    validation_rule: str = ""   # LLM 판단 — 계층 cascading 규칙·비고
    change_handler: str = ""    # 내부 — onChange leaf handler 이름 (trigger LLM 머지용)


@dataclass
class ScreenEvent:
    trigger: str = ""        # 버튼/이벤트 라벨
    event: str = ""          # onClick / onChange / componentDidMount ...
    backend_url: str = ""
    parent_handlers: List[str] = field(default_factory=list)  # this.props.X → 부모 함수 이름
    narrative: str = ""      # URL 무관 사이드이펙트 — "popup 열기 / 상태 갱신" 등
    source_offset: int = -1  # JSX 출현 순서 (정렬 secondary key, -1 = unknown)


@dataclass
class TableColumn:
    title: str = ""          # 필드설명 (한글 header, 예: "조직명")
    field: str = ""          # 필드명 (영문 dataIndex, 예: "org")
    width: str = ""          # 폭
    hide: bool = False       # 명시적 hidden=true / display:none
    # 화면정의서 9컬럼 표 양식
    data_type: str = ""      # 타입 (String / Number / Date)
    required: bool = False   # 필수여부
    attribute: str = ""      # 속성 (I/O/R/E/H 조합, 예: "O/R")
    ui_type: str = ""        # UI타입 (Text Field(Basic) / Dropdown / ...)
    description: str = ""    # 설명 (반환값)
    action: str = ""         # 동작 (예: "클릭시 [UI ID] 호출")
    # 분기 — JSX conditional ancestor (예: ``tab === 'A'``). 같은 화면에
    # 조건별로 다른 grid render 시 group key. 빈 값 = top-level (무조건).
    condition: str = ""
    # 길이 — 커스텀 cellEditor 의 maxLength 추출 (사용자 보고: ag-grid
    # ``cellEditor: 'ResetNumber'`` → ResetNumber.js 안 ``maxLength={10}``).
    length: str = ""


@dataclass
class ScreenLayout:
    file: str = ""           # React 파일 상대경로
    page_title: str = ""
    search_panel: List[ScreenField] = field(default_factory=list)
    # 입력 영역 — table 기반 입력 폼 (검색영역과 별도). panel_type='input'.
    input_panel: List[ScreenField] = field(default_factory=list)
    data_table_columns: List[TableColumn] = field(default_factory=list)
    edit_mode_fields: List[ScreenField] = field(default_factory=list)
    tabs: List[str] = field(default_factory=list)
    events: List[ScreenEvent] = field(default_factory=list)
    flowchart_mermaid: str = ""   # 사용자 액션 흐름 (Mermaid flowchart TB)
    summary: str = ""
    source: str = "llm"      # "llm" | "fallback"


# ── 캐시 ──────────────────────────────────────────────────────────────


def _cache_dir(base: str = "output/legacy_analysis/.screen_cache") -> Path:
    p = Path(base)
    p.mkdir(parents=True, exist_ok=True)
    return p


def _cache_key(file_content: str, url_map: Dict[str, List[str]]) -> str:
    payload = SCREEN_SCHEMA_VERSION + "\n" + file_content + "\n" + json.dumps(
        url_map, sort_keys=True
    )
    return hashlib.sha256(payload.encode("utf-8", errors="ignore")).hexdigest()


def _cache_get(key: str, enabled: bool) -> Optional[ScreenLayout]:
    if not enabled:
        return None
    fp = _cache_dir() / f"{key}.json"
    if not fp.exists():
        return None
    try:
        data = json.loads(fp.read_text(encoding="utf-8"))
        sp = [ScreenField(**f) for f in data.get("search_panel") or []]
        em = [ScreenField(**f) for f in data.get("edit_mode_fields") or []]
        ev = [ScreenEvent(**e) for e in data.get("events") or []]
        cols = [TableColumn(**c) if isinstance(c, dict) else TableColumn(title=str(c))
                for c in (data.get("data_table_columns") or [])]
        return ScreenLayout(
            file=data.get("file", ""),
            page_title=data.get("page_title", ""),
            search_panel=sp,
            data_table_columns=cols,
            edit_mode_fields=em,
            tabs=list(data.get("tabs") or []),
            events=ev,
            flowchart_mermaid=data.get("flowchart_mermaid", ""),
            summary=data.get("summary", ""),
            source=data.get("source", "llm"),
        )
    except Exception:
        return None


def _cache_put(key: str, layout: ScreenLayout, enabled: bool) -> None:
    if not enabled:
        return
    fp = _cache_dir() / f"{key}.json"
    try:
        fp.write_text(json.dumps(asdict(layout), ensure_ascii=False, indent=2),
                      encoding="utf-8")
    except Exception as e:
        logger.warning("screen cache 저장 실패 %s: %s", fp, e)


# ── LLM ──────────────────────────────────────────────────────────────


_SYSTEM_PROMPT = """당신은 React 화면 분석 전문가입니다. 주어진 React/JSX
파일을 분석해 화면 구조를 JSON 으로 추출하세요. 추측하지 말고 코드에서
명확히 읽히는 것만 채우세요. 발견 안 된 필드는 빈 배열/문자열로 두세요.

**중요**: events / backend_url 필드는 반환하지 마세요. 이벤트→백엔드 URL
매핑은 외부 정적 분석 결과를 사용하므로 LLM 응답에서 무시됩니다.
임의로 URL 을 추측하지 마세요.

**이미지가 첨부된 경우**: 첨부된 이미지는 **출력될 flowchart_mermaid
의 형태/스타일 sample** 입니다 (화면 스크린샷이 아닙니다). 해당 sample
의 노드 모양 / 분기 구조 / 화살표 방향 / 색상 / 그룹화 스타일을 그대로
참고해서 같은 시각 형태로 ``flowchart_mermaid`` 코드를 작성하세요.
화면 layout / search_panel / data_table_columns 추출과는 무관 — 그 필드들은
React 코드 텍스트만 보고 추출.

**검색 필드 라벨 추출 규칙** (한국 SI 흔한 패턴):
- 라벨은 JSX prop (``label`` / ``placeholder`` / ``title``) 으로 직접 들어
  오지 않고, 형제 element 의 text child 인 경우가 많음. 예::

      <div className="search-item">
        <span className="search-label required">Team</span>
        <Select defaultValue="Select 하세요." onChange={...}>
          <Option value="Y">Yes</Option>
          <Option value="N">No</Option>
        </Select>
      </div>

  → 라벨은 ``"Team"`` (형제 ``<span className="search-label">`` 의 text),
  컴포넌트는 ``"Select"``, default 는 ``"Select 하세요."``, options 는
  ``"Y, N"`` (Option 자식 value), events 는 ``"onChange"``, 라벨
  className 에 'required' 토큰 있으면 required=true.
- ``"Select 하세요."`` 는 default 값이지 라벨이 아님 — **혼동 금지**.
- prop 라벨이 비어있으면 형제 / ancestor (최대 5단계) 의 ``className`` 에
  'label' 이 포함된 span/div/label 의 text child 를 라벨로 사용.

**검색 영역 boundary**:
- 검색 패널 필드는 ``<section className="search-area">`` (또는
  ``search-form`` / ``filter-area`` / ``criteria-area``) 안의 input 컴포넌트
  **만** 포함. 그리드 inline filter / edit modal / 페이지네이션 등 다른
  영역의 input 은 search_panel 에서 제외.

**검색 패널 필드 9컬럼 양식** (data_table_columns 와 parallel):
- ``label`` — 라벨 텍스트 (위 추출 규칙 참조)
- ``input_data_type`` — **키보드 입력 필드일 때만** String/Number/Date 채움.
  Select/DatePicker/Checkbox/Radio 등 비-타이핑 필드는 빈 문자열.
- ``max_length`` — 키보드 입력 필드의 ``maxLength`` prop. 비-타이핑 필드는 빈.
- ``required`` — true/false (위 라벨 추출 규칙 참조)
- ``placeholder`` — ``placeholder`` prop 값 그대로 (있으면 default 보다 우선 표시)
- ``default`` — ``defaultValue`` / ``value`` literal (placeholder 없을 때만 표시)
- ``ui_type`` — Select(Single) / Select(Multi) / Text Field(Basic) /
  Text Field(Search Box) / DatePicker / Date Range / Checkbox / Radio Group /
  Number Field / Password / Text Area
- ``action`` — 단순 dropdown 이면 옵션 값 줄바꿈 (예: ``"전체\\nY\\nN"``).
  Cascading dependency (FAB → TEAM → SDPT 계층 — 부모 선택 시 자식 콤보 갱신)
  이 코드에서 보이면 그 동작 설명 (예: ``"FAB 선택 시 해당 TEAM 목록 조회"``).
- ``validation_rule`` — 검증 규칙 + 비고. 계층 의존성 hint (예: ``"FAB 미선택
  시 비활성"``) 또는 prop 기반 검증 (required/pattern/min/max) 요약.
  없으면 빈 문자열.

**DataTable 컬럼 추출 규칙** (화면정의서 9컬럼 표 양식):
- 컬럼 정의 prop 이름은 라이브러리별로 다름:
  ``columns`` (antd / react-table / generic) / ``columnDefs`` (ag-grid) /
  ``schema`` (RealGrid 등). React class state ``this.state.columnDefs`` 도
  처리 (state 객체 안 array 찾기).
- 컬럼 객체의 헤더 키 union: ``header`` / ``title`` / ``label`` /
  ``headerName`` (ag-grid). data 키 union: ``dataIndex`` / ``field`` /
  ``accessor`` / ``key`` / ``dataField``.
- 각 컬럼 한 행 — 9개 필드:
  - ``field`` — 영문 데이터 키 (예: "org")
  - ``title`` — 한글 필드설명 (예: "조직명")
  - ``data_type`` — String / Number / Date / Boolean 등
  - ``required`` — true/false (보통 그리드는 false, 검색폼은 true 많음)
  - ``attribute`` — I/O/R/E/H 조합 한 글자씩 슬래시 구분 (예: "O/R" /
    "O/E" / "H"). 기본은 "O/R" (Output + ReadOnly), ``editable: true`` →
    "O/E", ``hidden: true`` / ``hide: true`` → "H"
  - ``ui_type`` — "Text Field(Basic)" 기본, ``cellRenderer``/``cellEditor`` /
    ``type`` 단서로 "Dropdown" / "DatePicker" / "Number Field" / "Checkbox" /
    "Link/Button" 매핑
  - ``description`` — ``description`` / ``tooltipField`` / ``comment``
    prop 의 값 (반환값 / 데이터 설명). 없으면 빈 문자열
  - ``action`` — ``onCellClicked`` / ``onClick`` 등의 동작 요약 (예:
    "클릭시 [UI ID] 호출"). 없으면 빈 문자열
  - ``hide`` — bool (위 ``attribute`` 'H' 와 redundancy)

반환 schema (JSON):
{
  "page_title": "string — 화면 상단 제목",
  "search_panel": [
    {"label": "필드 라벨 (예: Team)",
     "component": "원본 JSX 컴포넌트 이름 (Select / DatePicker / Input / ...)",
     "default": "defaultValue / value (예: Select 하세요.)",
     "options": "드롭다운 자식 Option value 콤마 구분 (예: Y, N)",
     "events": "onChange / onClick / onBlur 등 (공백 구분)",
     "required": false,
     "placeholder": "placeholder attr 값 (예: Select 하세요)",
     "max_length": "maxLength prop 값 (keyboard 입력일 때만; 그 외 빈)",
     "input_data_type": "String | Number | Date (keyboard 입력일 때만; 그 외 빈)",
     "ui_type": "Select(Single) | Select(Multi) | Text Field(Basic) | Text Field(Search Box) | DatePicker | Date Range | Checkbox | Radio Group | Number Field | Password | Text Area",
     "action": "단순 dropdown 옵션 값 줄바꿈 (전체\\nY\\nN) 또는 cascading 동작 (FAB 선택 시 TEAM 갱신)",
     "validation_rule": "검증 규칙 / 계층 의존성 비고 (예: FAB 미선택 시 비활성) — 없으면 빈 문자열"}
  ],
  "data_table_columns": [
    {"field": "영문 데이터 키 / dataIndex (예: org)",
     "title": "한글 필드설명 (예: 조직명)",
     "data_type": "String | Number | Date | Boolean | ...",
     "required": false,
     "attribute": "O/R | O/E | H | I/O 등 슬래시 조합",
     "ui_type": "Text Field(Basic) | Dropdown | DatePicker | Number Field | Checkbox | Link/Button",
     "description": "반환값/데이터 설명 (예: 피평가자 조직 데이터)",
     "action": "동작 요약 (예: 클릭시 [UI ID] 호출(브라우저팝업)) — 없으면 빈 문자열",
     "width": "폭 표현 — 없으면 빈 문자열",
     "hide": false}
  ],
  "edit_mode_fields": [
    {"label": "...", "component": "...", "default": "...", "options": "..."}
  ],
  "tabs": ["탭1", "탭2", ...],
  "flowchart_mermaid": "사용자 액션 흐름 Mermaid flowchart TB 코드. 예시:\nflowchart TB\n    Start((화면 진입)) --> Init[초기 데이터 로드]\n    Init --> Display[그리드 표시]\n    Display --> Search{조회 클릭}\n    Search --> Update[그리드 갱신]\n    Display --> Detail{행 더블클릭}\n    Detail --> Popup[상세 popup 열림]\nMermaid v11 호환 규칙: (a) 라벨에 ``()`` / ``[]`` / ``/`` / ``:`` / ``=`` 등 특수문자 들어가면 반드시 double quote 로 감싸기. 예: Save[\"저장(POST)\"] / Cond{\"조회 N건\"} / Clear[\"초기화(targetKeys = [])\"]. (b) 노드 ID 는 영문/숫자/언더스코어만 (한글 ID 금지), 예약어 ``end`` / ``class`` / ``subgraph`` / ``style`` 사용 금지. (c) 라벨 안 nested ``[ ]`` / ``( )`` 두 단계 이상 nesting 은 피하고 평탄화 — nested 가 꼭 필요하면 라벨 전체 double quote. (d) 백엔드 URL 표시 X (events 표에 별도). 사용자 인터랙션 흐름만. (e) ``%%{init}%%`` 테마 directive 넣지 마. (f) 코드만 (```mermaid 펜스 X).",
  "summary": "1-2 줄 화면 설명"
}
"""


_RENDER_METHOD_RE = re.compile(r"^\s*render\s*\(\s*\)\s*\{", re.MULTILINE)
_FUNC_RETURN_JSX_RE = re.compile(r"return\s*\(\s*<", re.MULTILINE)
_IMPORT_LINE_RE = re.compile(r"^\s*import\s+", re.MULTILINE)


def _smart_slice(content: str, max_chars: int) -> str:
    """대용량 React 파일에서 LLM 분석에 필요한 부분만 추출.

    포함:
      1. imports 섹션 (라이브러리 imports — 컴포넌트 종류 단서)
      2. render() 메서드 본문 또는 functional component return JSX

    제외:
      - styled-components (긴 CSS template literals)
      - propTypes / defaultProps
      - render 와 무관한 helper functions
      - 주석 (regex slice 라 일부 남아있을 수 있음)

    파일이 max_chars 이하면 그대로 반환 (회귀 0).
    """
    if len(content) <= max_chars:
        return content

    parts: list[str] = []

    # 1. imports — 첫 import 부터 마지막 연속 import 라인까지
    imp_lines = []
    seen_import = False
    for line in content.split("\n"):
        s = line.lstrip()
        if s.startswith("import "):
            seen_import = True
            imp_lines.append(line)
        elif seen_import and (s.startswith("//") or not s):
            imp_lines.append(line)
        elif seen_import:
            break
    if imp_lines:
        parts.append("\n".join(imp_lines))
        parts.append("\n// ... (styled-components / helpers 생략) ...\n")

    # 2. render() 또는 functional return JSX block (brace walker)
    m = _RENDER_METHOD_RE.search(content)
    if m:
        start = m.start()
        i = content.index("{", m.start())
        depth = 0
        end = len(content)
        while i < len(content):
            c = content[i]
            if c == "{":
                depth += 1
            elif c == "}":
                depth -= 1
                if depth == 0:
                    end = i + 1
                    break
            i += 1
        parts.append(content[start:end])
    else:
        # functional: 첫 'return ( <' 부터 적당한 청크
        m = _FUNC_RETURN_JSX_RE.search(content)
        if m:
            chunk_size = max(2000, max_chars // 2)
            parts.append(content[m.start(): m.start() + chunk_size])

    result = "\n".join(parts)
    if len(result) > max_chars:
        result = result[:max_chars] + "\n... (smart slice truncated)"
    return result


def _build_user_prompt(file_rel: str, file_content: str,
                       url_map: Dict[str, List[str]],
                       max_chars: int,
                       *,
                       closure_markdown: Optional[str] = None) -> str:
    """LLM 사용자 프롬프트 빌더.

    ``closure_markdown`` 이 주어지면 `_smart_slice(file_content)` 대신
    closure 직렬화 결과 (entry + import 그래프 BFS 로 묶인 facts) 를
    소스 섹션으로 사용. ``--closure-llm`` 옵트인 경로.
    """
    if closure_markdown is not None:
        source_section = (
            "## React closure (import 그래프 BFS + popup 3 신호)\n"
            f"{closure_markdown}"
        )
    else:
        body = _smart_slice(file_content, max_chars)
        if len(file_content) > max_chars:
            body += (f"\n\n// (smart slice 적용: original {len(file_content)} chars "
                     f"→ {len(body)} chars 만 LLM 전달)")
        source_section = f"## React 소스\n```jsx\n{body}\n```"
    url_lines = []
    for handler, entry in sorted(url_map.items()):
        urls = entry.get("urls") if isinstance(entry, dict) else entry
        if urls:
            url_lines.append(f"  {handler} → {', '.join(sorted(set(urls)))}")
    url_block = "\n".join(url_lines) if url_lines else "  (없음)"
    return (
        f"파일: {file_rel}\n\n"
        f"## handler ↔ backend URL 매핑 (사전 정적 분석 결과)\n"
        f"{url_block}\n\n"
        f"{source_section}\n\n"
        "위 schema 의 JSON 만 반환하세요. 코드블록/설명 없이 raw JSON 만."
    )


_CLOSURE_IMPORT_WARNED = False


def _build_screen_closure(rel: str, abs_fp: str, frontend_dir: str,
                          patterns: Dict[str, Any],
                          max_depth: int, token_budget: int):
    """build_closure 한 번 호출 — markdown 직렬화 + 파서 양쪽이 같은 객체 공유.

    tree-sitter 미설치 / 빌드 실패 시 ``None``. 반환 타입은
    ``ScreenClosure`` 지만 import 가 lazy 라 forward-ref 로 둠.
    """
    global _CLOSURE_IMPORT_WARNED
    try:
        from .legacy_react_closure import build_closure
    except Exception as e:
        if not _CLOSURE_IMPORT_WARNED:
            logger.warning(
                "legacy_react_closure import 실패 (tree-sitter wheel 미설치?): %s "
                "— closure 의존 경로 비활성 (이후 화면들도 동일)", e
            )
            _CLOSURE_IMPORT_WARNED = True
        return None
    try:
        return build_closure(
            entry_file=abs_fp,
            repo_root=frontend_dir,
            patterns=patterns,
            max_depth=max_depth,
            token_budget=token_budget,
            verbose=False,
        )
    except Exception as e:
        logger.warning("closure build failed for %s: %s", rel, e)
        return None


def _serialize_closure_md(closure) -> Optional[str]:
    if closure is None:
        return None
    try:
        from .legacy_react_closure import serialize_for_llm
        return serialize_for_llm(closure)
    except Exception as e:
        logger.warning("closure serialize 실패: %s", e)
        return None


def _dump_screen_diagnostic(rel: str, closure) -> None:
    """search_panel / data_table_columns 둘 다 0 인 화면 — 사용자가 직접
    원인 추적할 수 있도록 closure 안 JSX 후보 1줄 dump.

    출력: ``[empty] <rel>: closure_files=N, table_candidates=[X(file:line) ...],
    input_candidates=[Y(file:line) ...]`` — 후보 tag 이름이 default
    패턴에 없으면 ``patterns.yaml`` 에 추가하면 됨.
    """
    try:
        from .legacy_react_ast import parse_file
    except Exception:
        return
    table_set = set(("Table", "DataTable", "Grid", "DataGrid", "AgGridReact",
                     "MaterialTable"))
    input_set = set(("input", "TextField", "TextInput", "Input", "Select",
                     "Dropdown", "DatePicker"))
    # closure 안 JSX tag 들 — table-like / input-like 후보만 count
    table_cands: list[str] = []
    input_cands: list[str] = []
    other_cap_cands: dict[str, int] = {}
    for f in closure.files:
        try:
            tree, source, _ = parse_file(f.abs_path)
        except Exception:
            continue
        if tree is None:
            continue
        from .legacy_react_ast import find_by_type, child_by_field, text_of
        for n in find_by_type(tree.root_node,
                              {"jsx_opening_element",
                               "jsx_self_closing_element"}):
            name_node = child_by_field(n, "name")
            if name_node is None:
                continue
            tag = text_of(name_node, source).strip()
            if not tag:
                continue
            line = n.start_point[0] + 1
            if tag in table_set:
                table_cands.append(f"{tag}({f.rel_path}:{line})")
            elif tag in input_set:
                input_cands.append(f"{tag}({f.rel_path}:{line})")
            elif tag and tag[0].isupper():
                # 대문자 시작 컴포넌트 — patterns.yaml 후보 (frequency 만)
                other_cap_cands[tag] = other_cap_cands.get(tag, 0) + 1
    others_top = sorted(other_cap_cands.items(), key=lambda x: -x[1])[:5]
    others_str = ", ".join(f"{n}({c})" for n, c in others_top) if others_top else "(none)"
    print(f"  [empty] {rel}: closure_files={len(closure.files)}, "
          f"tables={table_cands or '(none)'}, "
          f"inputs={(input_cands[:5] + ['...']) if len(input_cands) > 5 else (input_cands or ['(none)'])}, "
          f"other_custom_top5=[{others_str}]")


def _parser_fill_layout(layout: "ScreenLayout", closure,
                        patterns: Dict[str, Any]) -> tuple[int, int]:
    """파서 기반 search_panel / data_table_columns 덮어쓰기.

    closure 가 빌드되어 있으면 ``screen_spec.extractors`` 의 deterministic
    추출기로 폼 필드 / 그리드 컬럼을 다시 추출해서 LLM 결과를 덮어쓴다.
    events 와 동일 원칙: 파서가 ground truth.

    추출 결과 0 건이면 LLM 결과 유지 (회귀 회피).
    Returns ``(field_count, column_count)`` — 통계용.
    """
    if closure is None:
        return 0, 0
    try:
        from .screen_spec.extractors import (
            extract_form_fields,
            extract_grid_columns,
            extract_input_panel_fields,
        )
    except Exception as e:
        logger.warning("screen_spec extractors import 실패: %s", e)
        return 0, 0
    try:
        fields = extract_form_fields(closure, patterns)
        cols = extract_grid_columns(closure, patterns)
        input_fields = extract_input_panel_fields(closure, patterns)
    except Exception as e:
        logger.warning("파서 기반 화면 추출 실패: %s", e)
        return 0, 0
    if fields:
        # 파서 fill 전에 LLM 응답의 LLM-only 필드 (validation_rule / cascading
        # action) 를 라벨로 보존. 파서는 cascading 패턴을 못 추론하므로 LLM
        # 이 채워준 값을 잃지 않게 다시 머지. action 은 LLM 이 cascading
        # 설명을 채웠으면 (옵션 값 줄바꿈 list 가 아니면) 우선, 아니면 파서
        # 의 옵션 list 가 default.
        llm_extra: Dict[str, Dict[str, str]] = {}
        for sf in (layout.search_panel or []):
            key = (sf.label or "").strip()
            if not key:
                continue
            llm_extra[key] = {
                "validation_rule": (sf.validation_rule or ""),
                "action": (sf.action or ""),
            }

        def _merge_action(parser_act: str, llm_act: str) -> str:
            """parser_act (옵션 list 줄바꿈) vs llm_act (cascading 설명 등).

            LLM 이 의미있는 cascading 설명을 채웠으면 우선. LLM 값이 비었거나
            parser 와 동일한 옵션 list 이면 parser 유지.
            """
            llm_clean = (llm_act or "").strip()
            if not llm_clean:
                return parser_act or ""
            if llm_clean == (parser_act or "").strip():
                return parser_act or ""
            # LLM 값이 cascading 설명일 가능성 — 옵션 list 면 그대로, 아니면
            # parser 옵션 list 와 같이 보여줌 (LLM 설명이 우선 + 옵션 list 보강).
            if parser_act and parser_act.strip() != llm_clean:
                return llm_clean + "\n\n" + parser_act
            return llm_clean

        # component 는 원본 JSX tag 우선 (예: "Select", "DatePicker") —
        # field_type 의 소문자 분류 ("select"/"date") 대신 사용자가 보는
        # 화면에 실제 컴포넌트 이름이 그대로 표시되도록.
        layout.search_panel = [
            ScreenField(
                # name 은 식별자 (form key) 라 사용자 가시 라벨 아님 —
                # fallback 에서 제외. 라벨 없으면 빈 채로 두는 게 정답.
                label=(f.label or ""),
                name=(f.name or ""),
                component=(f.jsx_tag or f.field_type or ""),
                default=f.default or "",
                options=f.options or "",
                events=f.events or "",
                required=f.required,
                placeholder=f.placeholder or "",
                max_length=f.max_length or "",
                input_data_type=f.input_data_type or "",
                ui_type=f.ui_type or "",
                action=_merge_action(
                    f.action or "",
                    llm_extra.get((f.label or "").strip(), {}).get("action", ""),
                ),
                validation_rule=(
                    llm_extra.get((f.label or "").strip(), {})
                    .get("validation_rule", "")
                ),
                change_handler=getattr(f, "change_handler", "") or "",
            )
            for f in fields
        ]
    # input_panel — table 기반 입력 폼. search panel 과 같은 ScreenField
    # 구조 사용. parser 가 단독 제공 (LLM 머지 없음 — onSave 검증 자동
    # 추출이 deterministic 이라 충분).
    if input_fields:
        layout.input_panel = [
            ScreenField(
                label=(f.label or ""),
                name=(f.name or ""),
                component=(f.jsx_tag or f.field_type or ""),
                default=f.default or "",
                options=f.options or "",
                events=f.events or "",
                required=f.required,
                placeholder=f.placeholder or "",
                max_length=f.max_length or "",
                input_data_type=f.input_data_type or "",
                ui_type=f.ui_type or "",
                action=f.action or "",
                validation_rule=f.validation_rule or "",
                change_handler=getattr(f, "change_handler", "") or "",
            )
            for f in input_fields
        ]
    if cols:
        from .screen_spec.extractors import _compose_attribute
        # LLM 값 보존 — parser fill 전에 LLM 응답의 description / action /
        # ui_type 을 (title, field) 키로 캡처. 파서 후 빈 칸이면 LLM 값 재머지.
        # search_panel 과 같은 패턴.
        llm_grid_extra: Dict[str, Dict[str, str]] = {}
        for tc in (layout.data_table_columns or []):
            key = (tc.title or "").strip() or (tc.field or "").strip()
            if not key:
                continue
            llm_grid_extra[key.lower()] = {
                "description": (tc.description or ""),
                "action": (tc.action or ""),
                "ui_type": (tc.ui_type or ""),
            }

        def _llm_grid_pick(c, field: str) -> str:
            for key in (c.header or "", c.data_key or ""):
                k = (key or "").strip().lower()
                if k and k in llm_grid_extra:
                    return llm_grid_extra[k].get(field, "")
            return ""

        layout.data_table_columns = [
            TableColumn(
                title=c.header or "",
                field=c.data_key or "",
                width=c.width or "",
                hide=(not c.visible),
                data_type=(c.data_type or "").capitalize() or "String",
                required=c.required,
                attribute=_compose_attribute(c.visible, c.editable),
                ui_type=(c.ui_type or _llm_grid_pick(c, "ui_type")
                         or "Text Field(Basic)"),
                # description: parser (cellRenderer / onSave 검증 머지 결과) 우선,
                # 빈 칸이면 LLM 의 description 사용.
                description=(c.description or _llm_grid_pick(c, "description")),
                # action: parser 가 onCellClicked 등 인라인 처리. 빈 칸이면
                # LLM 의 action 설명 사용.
                action=(c.action or _llm_grid_pick(c, "action")),
                condition=getattr(c, "condition", "") or "",
                length=getattr(c, "length", "") or "",
            )
            for c in cols
        ]
    return len(fields), len(cols)


def _find_flowchart_sample(base_dir: str = "input") -> Optional[str]:
    """출력될 flowchart 의 형태/스타일 sample 이미지 — 모든 화면 공통 사용.

    사용자가 ``input/flowchart_sample.{png,jpg,jpeg,webp}`` 에 sample
    flowchart 이미지를 올리면 LLM 한테 첨부 → 같은 스타일로 flowchart_mermaid
    생성. 없으면 text-only 동작 (기존 그대로).
    """
    for name in ("flowchart_sample", "flowchart-sample"):
        for ext in ("png", "jpg", "jpeg", "webp"):
            p = os.path.join(base_dir, f"{name}.{ext}")
            if os.path.isfile(p):
                return p
    return None


# LLM prompt 크기 가드 — 너무 큰 prompt 는 LLM 게이트웨이가 silent fail
# 하거나 truncate 함. threshold 초과 시 source section 을 chunk 로 분할
# → N 회 호출 → JSON 응답 머지.
_LLM_PROMPT_CHUNK_THRESHOLD = 40000   # chars — 이걸 넘으면 chunk 모드
_LLM_PROMPT_CHUNK_SIZE = 30000        # 각 chunk 가 들어갈 안전한 크기


def _split_source_for_chunks(source: str, chunk_size: int) -> List[str]:
    """source 를 chunk_size 안에 맞게 분할.

    closure markdown 의 ``## File:`` 헤더 경계에서 자르려 시도 (파일
    통째로 한 chunk). 한 파일이 chunk_size 보다 크면 char 단위 분할
    (fallback). 분할이 필요 없으면 1-element list.
    """
    if len(source) <= chunk_size:
        return [source]
    # File boundary split — closure markdown 의 ``## File: ...`` 헤더 기준.
    # raw smart-slice 면 ``## File:`` 가 없을 텐데 그 경우 한 덩어리로
    # 재분할 (char-단위로 떨어짐).
    parts = re.split(r"(?=^## File:)", source, flags=re.MULTILINE)
    chunks: List[str] = []
    current = ""
    for p in parts:
        if not p:
            continue
        if len(current) + len(p) <= chunk_size:
            current += p
            continue
        if current:
            chunks.append(current)
            current = ""
        if len(p) <= chunk_size:
            current = p
        else:
            # 한 파일이 너무 큼 — char 단위 분할
            for i in range(0, len(p), chunk_size):
                chunks.append(p[i:i + chunk_size])
            current = ""
    if current:
        chunks.append(current)
    return chunks


def _merge_layout_dicts(dicts: List[Dict[str, Any]]) -> Dict[str, Any]:
    """여러 LLM 응답 dict 머지. chunk 응답 합치는 용도.

    - singles (page_title / summary / flowchart_mermaid): 첫 non-empty
    - lists (search_panel / input_panel / data_table_columns /
      edit_mode_fields / tabs): concat + dedupe (item 의 label/field/title
      키 기준)
    """
    merged: Dict[str, Any] = {
        "page_title": "",
        "search_panel": [],
        "input_panel": [],
        "data_table_columns": [],
        "edit_mode_fields": [],
        "tabs": [],
        "flowchart_mermaid": "",
        "summary": "",
    }

    def _item_key(item: Any, list_name: str) -> str:
        if not isinstance(item, dict):
            return str(item)
        # 필드별 dedup 키 — 같은 라벨/필드명이면 같은 항목
        if list_name == "data_table_columns":
            return str(item.get("field") or item.get("title") or item)
        return str(item.get("label") or item.get("name")
                   or item.get("id") or item)

    for d in dicts:
        if not isinstance(d, dict):
            continue
        for key in ("page_title", "flowchart_mermaid", "summary"):
            if not merged[key] and d.get(key):
                merged[key] = d[key]
        for key in ("search_panel", "input_panel", "data_table_columns",
                    "edit_mode_fields", "tabs"):
            existing_keys = {_item_key(it, key) for it in merged[key]}
            for item in (d.get(key) or []):
                k = _item_key(item, key)
                if k in existing_keys:
                    continue
                merged[key].append(item)
                existing_keys.add(k)
    return merged


def _call_llm_safe(prompt: str, config: Dict[str, Any],
                   label: str = "screen",
                   image_paths: Optional[list] = None) -> Optional[Dict[str, Any]]:
    try:
        from .legacy_pattern_discovery import _call_llm
    except Exception:
        return None
    try:
        raw = _call_llm(prompt, config or {}, label=label,
                        system_prompt=_SYSTEM_PROMPT,
                        image_paths=image_paths)
    except Exception as e:
        logger.warning("screen LLM 호출 실패: %s", e)
        return None
    if not raw:
        return None
    if isinstance(raw, dict):
        return raw
    if isinstance(raw, list) and raw and isinstance(raw[0], dict):
        return raw[0]
    return None


# ── Phase 3 — trigger 단위 LLM 분석 + 머지 ──


def _llm_analyze_triggers_for_screen(
    rel_path: str,
    content: str,
    url_map: Dict[str, Dict[str, Any]],
    config: Dict[str, Any],
    *,
    trigger_cache_dir: Optional[str] = None,
) -> Dict[str, dict]:
    """화면 1개의 모든 trigger 에 대해 LLM 분석 → ``{handler_name: analysis}``.

    url_map 은 ``_group_handlers_by_file`` 결과의 한 화면 entry —
    key 가 ``[event] label`` 형태라 event/handler/label 을 재추출 필요.
    각 trigger 마다 bundle 만들어 ``analyze_trigger_with_llm`` 호출.
    """
    try:
        from .legacy_react_api_scanner import (
            _build_call_regex, _DEFAULT_API_METHODS, _collect_url_constants,
            _collect_function_bodies, _collect_action_to_type,
            _collect_saga_urls_by_action_type,
            _collect_saga_fns_by_action_type, _collect_mdtp_action_map,
            _scan_dir, collect_event_handlers,
        )
        from .legacy_trigger_bundler import (
            build_trigger_bundle, analyze_trigger_with_llm,
        )
    except Exception as e:
        logger.warning("trigger LLM 모듈 import 실패: %s", e)
        return {}

    # 같은 file 안 trigger 모두 (url_map 기준). handler name 으로 dedup.
    handlers: Dict[str, dict] = {}
    for handler_label, entry in (url_map or {}).items():
        # handler_label 포맷: ``[onClick] 조회`` 또는 ``조회`` (event 없을 때)
        m = re.match(r"\[(?P<ev>[^\]]+)\]\s*(?P<label>.*)", handler_label)
        if m:
            ev = m.group("ev").strip()
            lbl = m.group("label").strip()
        else:
            ev, lbl = "", handler_label
        # event_marker 가 "onChange → parent.fnX" 같이 합성됐을 수 있어
        # 앞부분 event 만 사용.
        ev = ev.split()[0] if ev else ""
        handler = (entry.get("handler") if isinstance(entry, dict) else "") or ""
        # url_map 에 handler 가 따로 없으면 collect_event_handlers 결과로
        # label 매칭. 단순 fallback.
        if not handler:
            handler = lbl  # 가시 라벨로 대신 — LLM 이 jsx 보고 추론
        if handler not in handlers:
            handlers[handler] = {"event": ev, "label": lbl}

    if not handlers:
        return {}

    # collect_event_handlers 로 trigger entry 재추출 — body / source_offset
    # 같이 들고 옴 (bundle builder 에 필요).
    try:
        events_in_file = collect_event_handlers(content)
    except Exception:
        events_in_file = []
    by_handler = {(e.get("handler") or ""): e for e in events_in_file}

    # 인덱스 — 이미 caller (analyze_legacy) 가 만든 게 있을 수 있지만 단순
    # 화면 단위라 여기서 다시 빌드. 비싸지 않음.
    # ⚠ 비용: 같은 화면 frontend_dir 전체 walk → 다중 화면이면 caller
    # 가 한 번 만들어 넘기는 게 효율 (TODO Phase 후속 최적화).
    fe_dir = config.get("__frontend_dir") or ""
    if fe_dir and os.path.isdir(fe_dir):
        all_files = _scan_dir(fe_dir)
    else:
        all_files = []
    fn_index = _collect_function_bodies(all_files) if all_files else {}
    action_to_type = _collect_action_to_type(all_files) if all_files else {}
    call_re = _build_call_regex(list(_DEFAULT_API_METHODS))
    const_map = _collect_url_constants(all_files, []) if all_files else {}
    saga_urls = (_collect_saga_urls_by_action_type(
        all_files, fn_index, call_re, const_map, None) if all_files else {})
    saga_fns = (_collect_saga_fns_by_action_type(all_files)
                if all_files else {})
    mdtp_map = _collect_mdtp_action_map(all_files) if all_files else {}

    out: Dict[str, dict] = {}
    for handler_name, meta in handlers.items():
        if not handler_name:
            continue
        ev = by_handler.get(handler_name) or {
            "event": meta.get("event") or "",
            "handler": handler_name,
            "label": meta.get("label") or "",
            "body": "",
            "source_offset": -1,
        }
        try:
            bundle = build_trigger_bundle(
                ev, content, rel_path,
                fn_index=fn_index, mdtp_map=mdtp_map,
                action_to_type=action_to_type, saga_urls_by_type=saga_urls,
                saga_fns_by_type=saga_fns,
            )
        except Exception as e:
            logger.warning("trigger bundle 실패 (%s): %s", handler_name, e)
            continue
        analysis = analyze_trigger_with_llm(
            bundle, config, cache_dir=trigger_cache_dir, use_cache=True)
        if analysis:
            out[handler_name] = analysis
    return out


def _merge_trigger_llm_into_layout(layout: "ScreenLayout",
                                   trigger_analyses: Dict[str, dict]) -> None:
    """trigger LLM 분석 결과 → ScreenLayout 의 LLM-only 칸 머지.

    매칭 키:
    - search_panel field: change_handler / events (onChange 매칭)
    - events: trigger 의 handler/label
    - parser 가 채운 action / validation_rule 은 보존하면서 LLM 의
      더 풍부한 설명을 prepend.
    """
    if not trigger_analyses:
        return

    # search_panel — field.change_handler 정확 매칭 우선. 정확 매칭 없을
    # 때만 label substring fallback (사용자 보고: F/L 이 'fab' substring
    # 으로 FAB handler 결과를 잘못 가져옴 → 정확 매칭으로 회피).
    for f in layout.search_panel or []:
        matched = None
        # (a) 정확 매칭 — handler_name == f.change_handler
        ch = getattr(f, "change_handler", "") or ""
        if ch and ch in trigger_analyses:
            matched = trigger_analyses[ch]
        # (b) fallback — label 이 handler 이름에 substring 으로 들어있는 경우.
        # 단 너무 짧은 label (1-2 글자) 은 substring 매칭이 사고남 → skip.
        # (예: 'FL' 이 'handleFabChange' 안 'fl' substring 매칭하는 사고)
        if matched is None and f.label and len(f.label.strip()) >= 3:
            label_norm = re.sub(r"[\W_]+", "", f.label).lower()
            if label_norm:
                for handler_name, analysis in trigger_analyses.items():
                    if label_norm in handler_name.lower():
                        matched = analysis
                        break
        if not matched:
            continue
        # action — LLM 의 cascading 설명이 우선 + parser 옵션 list 보강.
        llm_action = (matched.get("action_description") or "").strip()
        if llm_action:
            if f.action and f.action.strip() != llm_action:
                f.action = llm_action + "\n\n" + f.action
            else:
                f.action = llm_action
        # validation_rule — LLM 만 채울 수 있는 칸. 기존 값 있으면 그대로
        # (parser cascading 검출 결과 보존).
        llm_val = (matched.get("validation_rule") or "").strip()
        if llm_val and not f.validation_rule:
            f.validation_rule = llm_val
        elif llm_val and llm_val not in (f.validation_rule or ""):
            f.validation_rule = (f.validation_rule + "\n" + llm_val).lstrip("\n")

    # events — handler 이름 매칭
    for ev in layout.events or []:
        # ScreenEvent 에 handler 필드가 따로 없으면 trigger 라벨로만 매칭.
        # narrative 칸에 business_summary 추가.
        for handler_name, analysis in trigger_analyses.items():
            label = ev.trigger or ""
            if label and (label in handler_name or handler_name in label):
                bs = (analysis.get("business_summary") or "").strip()
                if bs and bs not in (ev.narrative or ""):
                    ev.narrative = (ev.narrative + ", " + bs).lstrip(", ")
                break

    # required_fields — LLM 이 onSave / onSubmit 등 저장 흐름에서 필수
    # 검증을 받는 field 들을 식별. parser `_detect_save_validations` 가
    # 한국 SI 표준 패턴 (isNull / X===''||null / errorList.push 등) 은
    # 이미 잡지만, 커스텀 validation 라이브러리 / async 검증 / 비표준
    # 분기 같은 long-tail 은 LLM 이 보완. **parser 가 false 인 경우만**
    # LLM 응답으로 required=True 설정 (parser ground truth 우선).
    llm_required: set[str] = set()
    for analysis in trigger_analyses.values():
        for raw in (analysis.get("required_fields") or []):
            if not raw:
                continue
            # name 우선 매칭 — 정규화 (대소문자 / 특수문자 제거)
            llm_required.add(re.sub(r"[\W_]+", "", str(raw)).lower())
    if llm_required:
        for f in (layout.search_panel or []) + (layout.input_panel or []):
            if f.required:
                continue  # parser 가 이미 잡음 — 그대로
            keys = []
            if getattr(f, "name", ""):
                keys.append(re.sub(r"[\W_]+", "", f.name).lower())
            if f.label:
                keys.append(re.sub(r"[\W_]+", "", f.label).lower())
            if any(k and k in llm_required for k in keys):
                f.required = True


# ── Fallback (LLM 없을 때) — 정적 분석 결과만 가지고 events 만 채움 ──


def _fallback_layout(file_rel: str, url_map: Dict[str, Dict[str, Any]]) -> ScreenLayout:
    events: List[ScreenEvent] = []
    for handler, entry in url_map.items():
        # backward-compat: entry 가 plain list 면 기존 동작 유지.
        if isinstance(entry, list):
            urls = entry
            source_offset = -1
            narrative = ""
        else:
            urls = entry.get("urls") or []
            source_offset = int(entry.get("source_offset", -1))
            narrative = entry.get("narrative", "")
        m = re.match(r"\[(?P<ev>[^\]]+)\]\s*(?P<label>.*)", handler)
        if m:
            trigger = m.group("label").strip() or handler
            event_name = m.group("ev").strip()
        else:
            trigger, event_name = handler, ""
        if urls:
            for u in sorted(set(urls)):
                events.append(ScreenEvent(
                    trigger=trigger, event=event_name, backend_url=u,
                    narrative=narrative, source_offset=source_offset,
                ))
        else:
            # URL 무관 이벤트도 1 row — narrative 만 보여줌.
            events.append(ScreenEvent(
                trigger=trigger, event=event_name, backend_url="",
                narrative=narrative, source_offset=source_offset,
            ))
    return ScreenLayout(
        file=file_rel,
        page_title=os.path.splitext(os.path.basename(file_rel))[0],
        events=events,
        summary="LLM 없이 정적 분석 fallback — events 만 채움",
        source="fallback",
    )


def _parse_layout_dict(file_rel: str, data: Dict[str, Any]) -> ScreenLayout:
    def _fields(key: str) -> List[ScreenField]:
        out = []
        for f in data.get(key) or []:
            if isinstance(f, dict):
                out.append(ScreenField(
                    label=str(f.get("label", "")),
                    name=str(f.get("name", "") or f.get("id", "") or f.get("field", "")),
                    component=str(f.get("component", "")),
                    default=str(f.get("default", "")),
                    options=str(f.get("options", "")),
                    events=str(f.get("events", "")),
                    required=bool(f.get("required", False)),
                    placeholder=str(f.get("placeholder", "")),
                    max_length=str(f.get("max_length", "")),
                    input_data_type=str(f.get("input_data_type", "")),
                    ui_type=str(f.get("ui_type", "")),
                    action=str(f.get("action", "")),
                    validation_rule=str(f.get("validation_rule", "")),
                ))
        return out

    events: List[ScreenEvent] = []
    for e in data.get("events") or []:
        if isinstance(e, dict):
            events.append(ScreenEvent(
                trigger=str(e.get("trigger", "")),
                event=str(e.get("event", "")),
                backend_url=str(e.get("backend_url", "")),
            ))
    cols: List[TableColumn] = []
    for c in data.get("data_table_columns") or []:
        if isinstance(c, dict):
            cols.append(TableColumn(
                title=str(c.get("title", "")),
                field=str(c.get("field", "")),
                width=str(c.get("width", "")),
                hide=bool(c.get("hide", False)),
                data_type=str(c.get("data_type", "")),
                required=bool(c.get("required", False)),
                attribute=str(c.get("attribute", "")),
                ui_type=str(c.get("ui_type", "")),
                description=str(c.get("description", "")),
                action=str(c.get("action", "")),
                condition=str(c.get("condition", "")),
                length=str(c.get("length", "")),
            ))
        else:
            # 옛 형식 (단순 string) 호환
            cols.append(TableColumn(title=str(c)))
    return ScreenLayout(
        file=file_rel,
        page_title=str(data.get("page_title", "")),
        search_panel=_fields("search_panel"),
        data_table_columns=cols,
        edit_mode_fields=_fields("edit_mode_fields"),
        tabs=[str(t) for t in (data.get("tabs") or [])],
        events=events,
        flowchart_mermaid=str(data.get("flowchart_mermaid", "")),
        summary=str(data.get("summary", "")),
        source="llm",
    )


# ── 메인 추출 ─────────────────────────────────────────────────────────


def _group_handlers_by_file(handlers_by_url: Dict[str, List[Dict[str, Any]]]
                            ) -> Dict[str, Dict[str, Dict[str, Any]]]:
    """``{file: {handler_label: {urls, source_offset, narrative}}}`` 으로 변환.

    - parent_handlers 가 있으면 event_marker 에 ``→ parent.handleX`` 추가.
    - URL 무관 이벤트 (sentinel key ``""``) 도 emit — narrative 만 채워진
      엔트리로 들어감. ``--extract-screen-layout`` 화면 mockup 의 events
      테이블에 모든 버튼 표시 (URL 없어도 popup 열기 / 상태 갱신 narrative).
    """
    out: Dict[str, Dict[str, Dict[str, Any]]] = {}
    for url, ctx_list in (handlers_by_url or {}).items():
        for ctx in ctx_list or []:
            f = ctx.get("file") or ""
            if not f:
                continue
            handler = ctx.get("handler") or ""
            event_marker = ctx.get("event") or ""
            parents = ctx.get("parent_handlers") or []
            if parents:
                event_marker += " → " + ", ".join(f"parent.{p}" for p in parents)
            label = ctx.get("label") or ""
            tag = label or handler or "<inline>"
            full_handler = f"[{event_marker}] {tag}" if event_marker else tag
            entry = out.setdefault(f, {}).setdefault(full_handler, {
                "urls": [],
                "source_offset": ctx.get("source_offset", -1),
                "narrative": ctx.get("narrative", ""),
            })
            if url and url not in entry["urls"]:
                entry["urls"].append(url)
            # source_offset/narrative — 동일 handler 가 여러 url 로 emit 되면
            # 더 작은 (먼저 등장한) offset 으로 갱신, narrative 는 union.
            so = ctx.get("source_offset", -1)
            if so != -1 and (entry["source_offset"] == -1 or so < entry["source_offset"]):
                entry["source_offset"] = so
            n = ctx.get("narrative", "")
            if n and n not in entry["narrative"]:
                entry["narrative"] = (entry["narrative"] + ", " + n).lstrip(", ")
    return out


def extract_screen_layouts(
    frontend_dir: str,
    handlers_by_url: Dict[str, List[Dict[str, Any]]],
    patterns: Dict[str, Any],
    *,
    max_screens: int = 200,
    use_cache: bool = True,
    config: Optional[Dict[str, Any]] = None,
    closure_llm: bool = False,
    closure_max_depth: int = 3,
    closure_token_budget: int = 12000,
    llm_per_trigger: bool = False,
    trigger_cache_dir: Optional[str] = None,
) -> Dict[str, ScreenLayout]:
    """파일별로 한 번씩 LLM 호출 + 캐시. ``{file_rel: ScreenLayout}`` 반환.

    ``closure_llm=True`` (옵트인) 시 LLM input 을 raw JSX + smart_slice 대신
    AST 기반 closure markdown (import 그래프 BFS + popup 3 신호) 로 보강.
    tree-sitter 미설치/closure 빌드 실패 시 자동 smart_slice fallback.
    """
    cfg = dict(_DEFAULT_CONFIG)
    cfg.update((config or {}).get("screen_extraction") or {})
    max_chars = int(cfg.get("llm_max_chars", 6000))
    max_screens = max(1, min(max_screens, int(cfg.get("max_screens", 200))))

    by_file = _group_handlers_by_file(handlers_by_url)
    if not by_file:
        print("  screen layout: handler 컨텍스트 0건 — skip")
        return {}

    files = list(by_file.keys())
    if len(files) > max_screens:
        print(f"  screen layout: {len(files)} files > cap {max_screens} — truncate")
        files = files[:max_screens]

    print(f"  screen layout: {len(files)} React 화면 파일 분석 시작")

    out: Dict[str, ScreenLayout] = {}
    cache_hits = 0
    llm_calls = 0
    fallback_calls = 0
    closure_used = 0
    closure_failed = 0
    parser_fields_total = 0
    parser_grids_total = 0
    parser_screens = 0
    empty_screens = 0

    # 출력될 flowchart 형태 sample 이미지 — 한 번 lookup, 모든 화면 공통.
    sample_image = _find_flowchart_sample()
    if sample_image:
        print(f"  flowchart sample image: {sample_image}")

    if closure_llm:
        print(f"  closure_llm=ON (max_depth={closure_max_depth}, "
              f"token_budget={closure_token_budget})")

    for rel in files:
        url_map = by_file[rel]
        abs_fp = os.path.join(frontend_dir, rel)
        try:
            with open(abs_fp, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            # 주석 제거 — LLM 이 주석된 옛 코드를 보고 page_title /
            # search_panel 등을 잘못 추출하는 케이스 방지.
            from .legacy_react_api_scanner import _strip_comments
            content = _strip_comments(content)
        except Exception:
            content = ""

        # closure 한 번 빌드 — LLM markdown / 파서 기반 fill 양쪽이 공유.
        # tree-sitter 미설치면 None (회귀 0, LLM-only 경로 유지).
        closure_obj = _build_screen_closure(
            rel, abs_fp, frontend_dir, patterns or {},
            closure_max_depth, closure_token_budget,
        )
        closure_md: Optional[str] = None
        if closure_llm:
            closure_md = _serialize_closure_md(closure_obj)
            if closure_md:
                closure_used += 1
            else:
                closure_failed += 1

        # 캐시 키 — closure markdown 사용 시 raw content 대신 markdown 해시
        # (다른 입력 → 다른 LLM 결과). closure off 경로는 회귀 0 유지.
        cache_key = _cache_key(closure_md or content, url_map)
        cached = _cache_get(cache_key, use_cache)
        if cached:
            cached.file = rel
            out[rel] = cached
            cache_hits += 1
            continue

        prompt = _build_user_prompt(
            rel, content, url_map, max_chars,
            closure_markdown=closure_md,
        )
        # prompt 크기 가드 — threshold 초과면 source section 을 chunk 로
        # 분할 후 N 회 호출 + 응답 머지. 첫 chunk 의 응답을 base 로 하고
        # 나머지 chunk 의 lists 를 dedupe concat. LLM 게이트웨이의 hard
        # limit (보통 32K-128K context) 안전하게 통과.
        if len(prompt) > _LLM_PROMPT_CHUNK_THRESHOLD:
            source_section = closure_md if closure_md is not None else _smart_slice(content, max_chars)
            chunks = _split_source_for_chunks(
                source_section, _LLM_PROMPT_CHUNK_SIZE)
            print(f"  screen layout: prompt {len(prompt)} chars > "
                  f"{_LLM_PROMPT_CHUNK_THRESHOLD} — split {rel[:50]} "
                  f"into {len(chunks)} chunks")
            chunk_dicts: List[Dict[str, Any]] = []
            for i, chunk in enumerate(chunks, 1):
                chunk_prompt = _build_user_prompt(
                    rel,
                    chunk if closure_md is None else content,
                    url_map, max_chars,
                    closure_markdown=chunk if closure_md is not None else None,
                )
                chunk_data = _call_llm_safe(
                    chunk_prompt, config or {},
                    label=f"screen:{rel[:30]}:chunk{i}/{len(chunks)}",
                    # 이미지는 첫 chunk 만 (LLM 한테 같은 이미지 N번 보내지 않음)
                    image_paths=([sample_image] if (i == 1 and sample_image) else None),
                )
                if chunk_data:
                    chunk_dicts.append(chunk_data)
            data = _merge_layout_dicts(chunk_dicts) if chunk_dicts else None
        else:
            data = _call_llm_safe(
                prompt, config or {}, label=f"screen:{rel[:40]}",
                image_paths=[sample_image] if sample_image else None,
            )
        if data:
            layout = _parse_layout_dict(rel, data)
            llm_calls += 1
            # events 는 항상 정적 분석 결과로 덮어쓰기 — LLM 이 plausible 한
            # 환각 URL 만들어 진짜 호출 URL 가리는 케이스 차단. handlers_by_url
            # 는 collect_handler_contexts 가 JSX/saga 정적 분석으로 추출한
            # ground truth.
            layout.events = _fallback_layout(rel, url_map).events
        else:
            layout = _fallback_layout(rel, url_map)
            fallback_calls += 1

        # search_panel / data_table_columns 도 events 와 동일 원칙으로 파서가
        # ground truth. closure 가 빌드된 경우 (tree-sitter 설치 + 빌드 성공)
        # screen_spec.extractors 의 deterministic 추출기 결과로 덮어쓴다.
        # 자식 컴포넌트가 따로 import 된 분할 화면에서도 정확. 파서 0건이면
        # LLM 결과 유지.
        f_n, c_n = _parser_fill_layout(layout, closure_obj, patterns or {})
        if f_n or c_n:
            parser_screens += 1
            parser_fields_total += f_n
            parser_grids_total += c_n

        # Phase 3 — trigger 단위 LLM 분석 + 머지 (옵트인 ``--llm-per-trigger``).
        # 각 trigger 의 handler chain 을 한 덩어리로 LLM 에게 보내 동작 /
        # 유효성 / cascading 추출 → search_panel.action/validation_rule +
        # events.narrative 에 머지. parser facts (URL/setState) 는 ground
        # truth 그대로 유지.
        if llm_per_trigger:
            tlm = _llm_analyze_triggers_for_screen(
                rel, content, url_map, config or {},
                trigger_cache_dir=trigger_cache_dir,
            )
            if tlm:
                _merge_trigger_llm_into_layout(layout, tlm)

        # 진단: 파서 + LLM 모두 0 인 화면 — 사용자가 "그리드/조회영역 안
        # 나옴" 의심할 때 어디서 빠졌는지 보이도록 closure 안 JSX 후보를
        # 1줄 dump. closure 가 None 이면 (tree-sitter 미설치) skip.
        if (not layout.search_panel and not layout.data_table_columns
                and closure_obj is not None):
            empty_screens += 1
            _dump_screen_diagnostic(rel, closure_obj)

        out[rel] = layout
        _cache_put(cache_key, layout, use_cache)

    closure_stats = (
        f", closure_used={closure_used}, closure_failed={closure_failed}"
        if closure_llm else ""
    )
    parser_stats = (
        f", parser_screens={parser_screens}, "
        f"parser_fields={parser_fields_total}, parser_grids={parser_grids_total}"
        if parser_screens else ""
    )
    empty_stats = f", empty={empty_screens}" if empty_screens else ""
    print(f"  screen layout: cache_hits={cache_hits}, llm={llm_calls}, "
          f"fallback={fallback_calls}, total={len(out)}"
          f"{closure_stats}{parser_stats}{empty_stats}")
    return out


# ── HTML mockup 렌더 ──────────────────────────────────────────────────


_HTML_TEMPLATE = """<!doctype html>
<html lang="ko"><head><meta charset="utf-8">
<title>{title}</title>
<style>
  body {{ font-family: -apple-system, "맑은 고딕", sans-serif; margin: 0;
         background: #f4f5f7; color: #222; }}
  header {{ background: #2c3e50; color: #fff; padding: 12px 20px;
            font-size: 18px; font-weight: 600; }}
  main {{ max-width: 1100px; margin: 16px auto; padding: 0 16px; }}
  section {{ background: #fff; border-radius: 4px; padding: 16px;
             margin-bottom: 12px; box-shadow: 0 1px 2px rgba(0,0,0,0.05); }}
  section h2 {{ font-size: 14px; color: #555; margin: 0 0 12px;
                border-bottom: 1px solid #eee; padding-bottom: 6px; }}
  .field {{ display: inline-block; margin: 4px 12px 4px 0; }}
  .field label {{ font-size: 12px; color: #666; display: block; margin-bottom: 2px; }}
  .field input, .field select {{ border: 1px solid #ccc; padding: 4px 8px;
            background: #fafafa; min-width: 140px; }}
  .field .note {{ font-size: 11px; color: #999; margin-left: 4px; }}
  table {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
  table.grid-def th {{ font-size: 12px; background: #2c3e50; color: #fff; }}
  table.grid-def td {{ vertical-align: top; }}
  table.grid-def td.no {{ text-align: center; color: #777; width: 36px; }}
  table.grid-def code {{ font-size: 12px; color: #2c3e50; }}
  .req-badge {{ display: inline-block; font-size: 10px; padding: 1px 6px;
              background: #c0392b; color: #fff; border-radius: 2px;
              margin-left: 6px; vertical-align: middle; }}
  th, td {{ border: 1px solid #ddd; padding: 6px 8px; text-align: left; }}
  th {{ background: #f0f0f0; }}
  td.placeholder {{ color: #bbb; font-style: italic; }}
  .tab-bar {{ display: flex; border-bottom: 2px solid #2c3e50; }}
  .tab {{ padding: 8px 16px; cursor: pointer; background: #ecf0f1;
         border: 1px solid #ddd; border-bottom: 0; margin-right: 2px; }}
  .tab.active {{ background: #2c3e50; color: #fff; }}
  .events {{ font-size: 12px; }}
  .events table th {{ background: #fffbe5; }}
  pre.mermaid {{ background: #f8f9fa; border: 1px solid #e0e0e0;
                  padding: 12px; font-size: 12px; overflow-x: auto;
                  font-family: Consolas, "Courier New", monospace; }}
  .meta {{ font-size: 11px; color: #999; margin-top: 16px;
           text-align: right; }}
  .empty {{ color: #aaa; font-style: italic; }}
</style>
<script src="https://cdn.jsdelivr.net/npm/mermaid/dist/mermaid.min.js"></script>
<script>
// Mermaid 11.x 는 strict + parse 실패 시 throw 대신 에러 SVG 를 그려 반환.
// → suppressErrorRendering: true 로 throw 강제 + parse() 로 사전 검증 → catch
// 가 raw 텍스트 + 에러 메시지를 노출 (사용자 환경 단방향이라 디버깅 단서 필수).
if (window.mermaid) {{
  mermaid.initialize({{
    startOnLoad: false,
    securityLevel: 'loose',
    suppressErrorRendering: true
  }});
  window.addEventListener('DOMContentLoaded', async function () {{
    var blocks = document.querySelectorAll('pre.mermaid');
    var ver = '?';
    try {{
      if (typeof mermaid.version === 'function') ver = mermaid.version();
      else if (typeof mermaid.version === 'string') ver = mermaid.version;
      else if (mermaid.mermaidAPI && mermaid.mermaidAPI.version) {{
        ver = (typeof mermaid.mermaidAPI.version === 'function')
              ? mermaid.mermaidAPI.version() : mermaid.mermaidAPI.version;
      }}
    }} catch (e) {{}}
    if (ver === '?') {{
      // CDN URL 에서 version 추출 시도
      var scripts = document.querySelectorAll('script[src*="mermaid"]');
      for (var s = 0; s < scripts.length; s++) {{
        var mm = scripts[s].src.match(/mermaid[^@]*@?([\\d.]+)/);
        if (mm) {{ ver = mm[1]; break; }}
      }}
    }}
    for (var i = 0; i < blocks.length; i++) {{
      var el = blocks[i];
      var src = el.textContent;
      var lastErr = null;
      try {{
        await mermaid.parse(src);   // parse 가 명확한 syntax error throw
        var id = 'm' + i + '_' + Math.random().toString(36).slice(2);
        var out = await mermaid.render(id, src);
        // render 가 어쩌다 success 하지만 SVG 안에 에러를 그린 케이스 방어 —
        // SVG 가 'aria-roledescription="error"' 또는 'mermaid-error' 클래스
        // 포함하면 실패로 간주.
        if (/aria-roledescription="error"|class="error-/.test(out.svg)) {{
          throw new Error('mermaid render returned error SVG');
        }}
        el.innerHTML = out.svg;
        continue;
      }} catch (e) {{
        lastErr = e;
      }}
      var msg = (lastErr && (lastErr.message || lastErr.str)) || String(lastErr);
      el.innerHTML =
        '<div style="color:#c0392b;font-weight:600;margin-bottom:6px;">' +
        'Mermaid parse error (v' + ver + '): ' + msg + '</div>' +
        '<pre style="background:#fff;color:#333;border:1px dashed #c0392b;' +
        'padding:8px;white-space:pre-wrap;font-family:Consolas,monospace;">' +
        src.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;') +
        '</pre>';
    }}
  }});
}}
</script>
</head><body>
<header>{title_html}</header>
<main>
{summary_block}
{search_block}
{input_block}
{tab_block}
{table_block}
{flowchart_block}
{edit_block}
{events_block}
<div class="meta">file: {file_rel} · source: {source}</div>
</main></body></html>
"""


def _esc(s: str) -> str:
    return html.escape(s or "", quote=True)


def _render_field_list(fields: List[ScreenField]) -> str:
    """필드 목록 → 텍스트 bullet 리스트. 항목당::

        <strong>라벨</strong> [필수]
          - 사용 컴포넌트: Select
          - Default: ...
          - 리스트: Y, N
          - 이벤트: onChange
    """
    items = []
    for f in fields:
        sub = []
        if f.component:
            sub.append(f"<li>사용 컴포넌트: {_esc(f.component)}</li>")
        if f.default:
            sub.append(f"<li>Default: {_esc(f.default)}</li>")
        if f.options:
            sub.append(f"<li>리스트: {_esc(f.options)}</li>")
        if f.events:
            sub.append(f"<li>이벤트: {_esc(f.events)}</li>")
        sub_html = f"<ul>{''.join(sub)}</ul>" if sub else ""
        req_badge = (" <span class='req-badge'>필수</span>" if f.required else "")
        items.append(
            f"<li><strong>{_esc(f.label or '(no label)')}</strong>"
            f"{req_badge}{sub_html}</li>"
        )
    return f"<ol>{''.join(items)}</ol>" if items else ""


def _render_search(fields: List[ScreenField]) -> str:
    if not fields:
        return ""
    return ("<section><h2>Search Panel — 조회 조건 영역</h2>"
            + _render_field_list(fields)
            + _render_search_table(fields)
            + "</section>")


def _render_input_panel(fields: List[ScreenField]) -> str:
    """입력 영역 (table 기반 입력 폼) 정의서 — 검색영역과 같은 컬럼 양식.

    컬럼: No / 필드(영문) / 라벨 / 타입 / 길이 / 필수 / 기본값 / 유효성 규칙 및 비고 /
    UI 타입 / 동작.

    onSave 의 isNull 검증은 ``required`` 로, isNumber/isNegative 등 그 외
    검증은 ``validation_rule`` 로 parser 가 채움.
    """
    if not fields:
        return ""
    rows = []
    for i, f in enumerate(fields, start=1):
        display_default = f.placeholder or f.default or ""
        action_html = _esc(f.action or "").replace("\n", "<br/>")
        field_en = _esc(getattr(f, "name", "") or "")
        rows.append(
            "<tr>"
            f"<td class='no'>{i}</td>"
            f"<td><code>{field_en}</code></td>"
            f"<td>{_esc(f.label or '')}</td>"
            f"<td>{_esc(f.input_data_type or '')}</td>"
            f"<td>{_esc(f.max_length or '')}</td>"
            f"<td>{'필수' if f.required else '선택'}</td>"
            f"<td>{_esc(display_default)}</td>"
            f"<td>{_esc(f.validation_rule or '')}</td>"
            f"<td>{_esc(f.ui_type or '')}</td>"
            f"<td>{action_html}</td>"
            "</tr>"
        )
    return (
        "<section><h2>Input Panel — 입력 영역 정의서</h2>"
        "<table class='input-spec'><thead><tr>"
        "<th>No</th><th>필드(영문)</th><th>라벨</th><th>타입</th><th>길이</th>"
        "<th>필수</th><th>기본값</th><th>유효성 규칙 및 비고</th>"
        "<th>UI 타입</th><th>동작</th>"
        "</tr></thead><tbody>"
        + "".join(rows)
        + "</tbody></table></section>"
    )


def _render_search_table(fields: List[ScreenField]) -> str:
    """검색 패널 9컬럼 화면정의서 표 (grid 와 parallel).

    컬럼: No / 필드(영문) / 라벨 / 타입 / 길이 / 필수 / 기본값 / 유효성 규칙 및 비고 /
    UI 타입 / 동작. 기본값은 placeholder 우선 (UI 가시값) → default fallback.
    동작은 단순 dropdown 이면 옵션 값 줄바꿈, LLM 이 cascading 동작
    채웠으면 그 값 우선.

    필드(영문) — input element 의 id 우선, 없으면 name (ScreenField.name 에
    이미 그 순서로 추출돼있음).
    """
    if not fields:
        return ""
    rows = []
    for i, f in enumerate(fields, start=1):
        display_default = f.placeholder or f.default or ""
        validation_display = f.validation_rule or ""
        action_html = _esc(f.action or "").replace("\n", "<br/>")
        field_en = _esc(getattr(f, "name", "") or "")
        rows.append(
            "<tr>"
            f"<td class='no'>{i}</td>"
            f"<td><code>{field_en}</code></td>"
            f"<td>{_esc(f.label or '')}</td>"
            f"<td>{_esc(f.input_data_type or '')}</td>"
            f"<td>{_esc(f.max_length or '')}</td>"
            f"<td>{'필수' if f.required else '선택'}</td>"
            f"<td>{_esc(display_default)}</td>"
            f"<td>{_esc(validation_display)}</td>"
            f"<td>{_esc(f.ui_type or '')}</td>"
            f"<td>{action_html}</td>"
            "</tr>"
        )
    return (
        "<h3 style='margin-top:12px'>검색 영역 정의서</h3>"
        "<table class='search-spec'><thead><tr>"
        "<th>No</th><th>필드(영문)</th><th>라벨</th><th>타입</th><th>길이</th>"
        "<th>필수</th><th>기본값</th><th>유효성 규칙 및 비고</th>"
        "<th>UI 타입</th><th>동작</th>"
        "</tr></thead><tbody>"
        + "".join(rows)
        + "</tbody></table>"
    )


def _render_table(cols: List[TableColumn]) -> str:
    """화면정의서 표 양식 — NO / 필드명(영문) / 필드설명 / 타입 / 필수여부 /
    속성 / UI타입 / 설명 / 동작 9컬럼.

    grid 의 ``condition`` 이 다르면 condition 별로 sub-section 분리
    (``{tab === 'A' && <Grid/>}`` 같이 분기 render 되는 grid). 같은
    condition 안 컬럼들은 한 표로 묶음.

    hide 컬럼은 attribute='H' 로 표시하되 같은 표에 emit (별도 분리 X).
    """
    if not cols:
        return ""
    # condition 별 group — 순서 보존 (먼저 등장한 condition 부터).
    from collections import OrderedDict as _OD
    groups: "_OD[str, List[TableColumn]]" = _OD()
    for c in cols:
        cond = getattr(c, "condition", "") or ""
        groups.setdefault(cond, []).append(c)
    if len(groups) == 1 and "" in groups:
        # condition 없는 단일 그룹 — 기존 동작 (sub-section 분리 X)
        return _render_table_inner(groups[""], heading=None)
    sections = []
    for cond, group_cols in groups.items():
        heading = f"분기: {cond}" if cond else "분기 없음 (top-level)"
        sections.append(_render_table_inner(group_cols, heading=heading))
    return "".join(sections)


def _render_table_inner(cols: List[TableColumn], heading: str | None) -> str:
    """내부 헬퍼 — 단일 grid (또는 단일 condition group) 의 표 렌더.

    컬럼 순서 (사용자 명시):
      NO / 필드명(영문) / 헤더명 / 타입 / 길이 / 필수여부 / 속성 / 필드설명
      / UI타입 / 동작
    """
    headers = ["NO", "필드명(영문)", "헤더명", "타입", "길이", "필수여부",
               "속성", "필드설명", "UI타입", "동작"]
    head_row = "".join(f"<th>{h}</th>" for h in headers)
    body_rows = []
    for i, c in enumerate(cols, start=1):
        attribute = c.attribute or ("H" if c.hide else "O/R")
        ui = c.ui_type or "Text Field(Basic)"
        dtype = (c.data_type or "String").capitalize()
        required = "필수" if c.required else "선택"
        length = getattr(c, "length", "") or ""
        row = (
            f"<tr>"
            f"<td class='no'>{i}</td>"
            f"<td><code>{_esc(c.field) or '<em>(no field)</em>'}</code></td>"
            f"<td>{_esc(c.title)}</td>"
            f"<td>{_esc(dtype)}</td>"
            f"<td>{_esc(length)}</td>"
            f"<td>{_esc(required)}</td>"
            f"<td>{_esc(attribute)}</td>"
            f"<td>{_esc(c.description)}</td>"
            f"<td>{_esc(ui)}</td>"
            f"<td>{_esc(c.action)}</td>"
            f"</tr>"
        )
        body_rows.append(row)
    title = "DataTable (그리드 정의)"
    if heading:
        title = f"DataTable — {heading}"
    return (
        f"<section><h2>{_esc(title)}</h2>"
        f"<table class='grid-def'><thead><tr>{head_row}</tr></thead>"
        f"<tbody>{''.join(body_rows)}</tbody></table>"
        "</section>"
    )


def _render_edit(fields: List[ScreenField]) -> str:
    if not fields:
        return ""
    return ("<section><h2>Edit Mode — 편집 영역</h2>"
            + _render_field_list(fields) + "</section>")


def _render_tabs(tabs: List[str]) -> str:
    if not tabs:
        return ""
    items = []
    for i, t in enumerate(tabs):
        cls = "tab active" if i == 0 else "tab"
        items.append(f"<div class='{cls}'>{_esc(t)}</div>")
    return f"<section><h2>Tabs</h2><div class='tab-bar'>{''.join(items)}</div></section>"


def _event_sort_rank(event: str) -> int:
    """이벤트 정렬 우선순위 (사용자 요청 — 화면오픈 → onChange → onClick → 그 외).

    낮을수록 먼저 표시.
    """
    e = (event or "").lower()
    if any(s in e for s in ("mount", "useeffect", "didmount", "willmount")):
        return 0   # 화면 최초 오픈
    if "didupdate" in e:
        return 1   # 업데이트
    if "change" in e:
        return 2   # onChange / onValueChange
    if "submit" in e:
        return 3   # form submit
    if "click" in e:
        return 4   # 버튼 클릭
    return 5       # 그 외 (onBlur, onFocus, onKeyDown, ...)


def _render_events(events: List[ScreenEvent]) -> str:
    """Trigger + Event 별로 그룹화 후 한 row 에 backend URL + narrative.

    정렬: lifecycle (mount) → onChange → onClick → 나머지 (rank 우선),
    같은 rank 안에서는 **JSX 출현 순서** (source_offset) → trigger 라벨.
    URL 없는 버튼 (popup 호출 / 상태 갱신 등) 도 narrative 열에 표시.
    """
    if not events:
        return ""
    # (trigger, event) → {urls, source_offset, narrative}
    grouped: dict[tuple[str, str], dict[str, Any]] = {}
    for e in events:
        key = (e.trigger or "<inline>", e.event or "")
        rec = grouped.setdefault(key, {
            "urls": [], "source_offset": e.source_offset, "narrative": e.narrative,
        })
        if e.backend_url and e.backend_url not in rec["urls"]:
            rec["urls"].append(e.backend_url)
        if e.source_offset != -1 and (rec["source_offset"] == -1
                                      or e.source_offset < rec["source_offset"]):
            rec["source_offset"] = e.source_offset
        if e.narrative and e.narrative not in rec["narrative"]:
            rec["narrative"] = (rec["narrative"] + ", " + e.narrative).lstrip(", ")
    # 정렬: rank → source_offset (없으면 매우 큼) → trigger
    def _sort_key(item):
        (trig, ev), rec = item
        so = rec["source_offset"]
        so_key = so if so != -1 else 1_000_000
        return (_event_sort_rank(ev), so_key, trig)
    rows = []
    has_narrative = any(rec["narrative"] for rec in grouped.values())
    for (trigger, event), rec in sorted(grouped.items(), key=_sort_key):
        url_html = "<br>".join(f"<code>{_esc(u)}</code>" for u in rec["urls"]) or "—"
        narr_html = _esc(rec["narrative"]) or "—"
        if has_narrative:
            rows.append(
                f"<tr><td>{_esc(trigger)}</td><td>{_esc(event)}</td>"
                f"<td>{url_html}</td><td>{narr_html}</td></tr>"
            )
        else:
            rows.append(
                f"<tr><td>{_esc(trigger)}</td><td>{_esc(event)}</td>"
                f"<td>{url_html}</td></tr>"
            )
    head_extra = "<th>설명</th>" if has_narrative else ""
    return (f"<section class='events'><h2>이벤트 → 백엔드 URL</h2>"
            f"<table><thead><tr><th>Trigger</th><th>Event</th>"
            f"<th>Backend URL</th>{head_extra}</tr></thead>"
            f"<tbody>{''.join(rows)}</tbody></table></section>")


_MERMAID_DIRECTIVE_RE = re.compile(r"%%\{.*?\}%%", re.DOTALL)
# 노드 시작 — `ID[` / `ID(` / `ID{` 등.  ID 는 영문/숫자/언더스코어/하이픈.
_MERMAID_NODE_OPEN_RE = re.compile(r"\b([A-Za-z_][\w\-]*)([\[\(\{])")
# Arrow / link — 라벨 종료 후보 (라벨 검색 범위 boundary).
_MERMAID_ARROW_RE = re.compile(r"-{2,}>|={2,}>|-\.+->|-{2,}|\.{2,}|\|")
# Mermaid 11.x reserved keywords — 노드 ID 로 쓰면 parse 실패.
_MERMAID_RESERVED_IDS = {"end", "class", "subgraph", "style", "default", "linkStyle"}
# Label 안 risky chars — 있으면 ``"..."`` quote + HTML entity escape.
_MERMAID_RISKY = set("()[]{}:;,/<>=")
_HTML_ENTITY_BRACKETS = str.maketrans({
    "[": "&#91;", "]": "&#93;",
    "(": "&#40;", ")": "&#41;",
    "{": "&#123;", "}": "&#125;",
})


def _sanitize_node_labels(text: str) -> str:
    """nested bracket-aware 라벨 quoting.

    노드 ID 직후 열린 ``[`` / ``(`` / ``{`` 부터 **라인 끝 또는 다음 arrow
    직전까지의 영역에서 마지막 close bracket** 까지를 라벨로 간주.
    ``A[ClearSelect(... [])]`` 처럼 라벨 안에 nested brackets 가 있어도
    boundary 정확히 잡음 (regex lazy 매칭 + inner bracket 차단 문제 해소).
    라벨에 risky char 있으면 inner brackets 를 HTML entity 로 escape 후
    ``"..."`` 로 감쌈 — Mermaid v10+ 가 entity 인식.
    """
    out: list[str] = []
    pos = 0
    n = len(text)
    while pos < n:
        m = _MERMAID_NODE_OPEN_RE.search(text, pos)
        if not m:
            out.append(text[pos:])
            break
        out.append(text[pos:m.start()])
        node_id = m.group(1)
        open_ch = m.group(2)
        open_pos = m.end() - 1
        # double opener?
        open_count = 1
        if open_pos + 1 < n and text[open_pos + 1] == open_ch:
            open_count = 2
        close_ch = {"[": "]", "(": ")", "{": "}"}[open_ch]
        section_start = open_pos + open_count
        # boundary = 라인 끝 OR 다음 arrow
        end_of_line = text.find("\n", section_start)
        if end_of_line == -1:
            end_of_line = n
        arrow_m = _MERMAID_ARROW_RE.search(text, section_start, end_of_line)
        section_end = arrow_m.start() if arrow_m else end_of_line
        section = text[section_start:section_end]
        if open_count == 2:
            cc = close_ch * 2
            last_idx = section.rfind(cc)
            if last_idx == -1:
                out.append(text[m.start():section_start])
                pos = section_start
                continue
            label = section[:last_idx]
            end_abs = section_start + last_idx + 2
        else:
            last_idx = section.rfind(close_ch)
            if last_idx == -1:
                out.append(text[m.start():section_start])
                pos = section_start
                continue
            label = section[:last_idx]
            end_abs = section_start + last_idx + 1

        stripped = label.strip()
        already_quoted = stripped.startswith('"') and stripped.endswith('"')
        if already_quoted or not any(c in label for c in _MERMAID_RISKY):
            out.append(text[m.start():end_abs])
        else:
            safe = label.replace('"', "'").translate(_HTML_ENTITY_BRACKETS)
            out.append(node_id + open_ch * open_count + '"' + safe + '"'
                       + close_ch * open_count)
        pos = end_abs
    return "".join(out)


def _sanitize_mermaid_flowchart(code: str) -> str:
    """LLM 이 생성한 Mermaid 코드를 11.x parser 친화적으로 정리.

    - ```mermaid 펜스 제거
    - %%{init: ...}%% 테마 directive 제거 (LLM 이 잘못 inject 하는 경우 방어)
    - 노드 라벨 quote + HTML entity escape (nested brackets 대응)
    - 예약어 노드 ID 충돌 회피 (``end`` → ``end_`` 등)
    - flowchart directive 누락 시 ``flowchart TB`` 자동 prepend
    """
    if not code:
        return ""
    s = code.strip()
    if s.startswith("```"):
        s = s.strip("`")
        first_nl = s.find("\n")
        if first_nl > 0 and "mermaid" in s[:first_nl].lower():
            s = s[first_nl + 1:]
        s = s.rstrip("` \n")
    s = _MERMAID_DIRECTIVE_RE.sub("", s)
    s = _sanitize_node_labels(s)

    # 예약어 노드 ID 치환 — 라벨 시작 / 화살표 양옆 / 줄끝 직전 모두 cover.
    def _rename_reserved(text: str) -> str:
        for kw in _MERMAID_RESERVED_IDS:
            text = re.sub(
                rf"(?<![\w]){kw}(?=\s*[\[\(\{{]|\s*-->|\s*--|\s*$)",
                f"{kw}_", text, flags=re.MULTILINE,
            )
        return text
    s = _rename_reserved(s)

    # flowchart 헤더 보장
    lines = [ln for ln in s.splitlines() if ln.strip()]
    if not lines:
        return ""
    first = lines[0].strip().lower()
    if not (first.startswith("flowchart") or first.startswith("graph")):
        lines.insert(0, "flowchart TB")
    return "\n".join(lines)


def _render_flowchart(mermaid_code: str) -> str:
    """사용자 액션 흐름 — Mermaid flowchart. <pre class='mermaid'> 안에
    sanitize 한 코드. 페이지 JS 가 mermaid.render() 로 그리며 parse 실패 시
    raw 코드 + 에러를 그대로 노출 (디버깅용).
    LLM 이 ``flowchart_mermaid`` 빈 값 반환하면 섹션 자체 미생성.
    """
    code = _sanitize_mermaid_flowchart(mermaid_code or "")
    if not code:
        return ""
    return ("<section><h2>화면 흐름 (사용자 액션)</h2>"
            f"<pre class='mermaid'>{_esc(code)}</pre></section>")


def render_screen_html(layout: ScreenLayout) -> str:
    title = layout.page_title or os.path.basename(layout.file) or "Screen"
    summary_block = (f"<section><h2>요약</h2><div>{_esc(layout.summary)}</div></section>"
                     if layout.summary else "")
    return _HTML_TEMPLATE.format(
        title=_esc(title),
        title_html=_esc(title),
        summary_block=summary_block,
        search_block=_render_search(layout.search_panel),
        input_block=_render_input_panel(layout.input_panel or []),
        tab_block=_render_tabs(layout.tabs),
        table_block=_render_table(layout.data_table_columns),
        flowchart_block=_render_flowchart(layout.flowchart_mermaid),
        edit_block=_render_edit(layout.edit_mode_fields),
        events_block=_render_events(layout.events),
        file_rel=_esc(layout.file),
        source=_esc(layout.source),
    )


def _find_mmdc_executable() -> Optional[str]:
    """``mmdc`` (mermaid-cli) executable 위치 1 회 lookup. PATH 에 없으면 None.

    Windows 면 ``mmdc.cmd`` 도 같이 검색.
    """
    return shutil.which("mmdc") or shutil.which("mmdc.cmd")


def _render_mmd_to_format(mmdc_path: str, mmd_path: str, ext: str,
                           timeout: int = 30, width: int = 1920) -> Optional[str]:
    """mmdc 로 .mmd → .{ext} 변환 (ext: svg / png / pdf). 성공 시 출력
    경로, 실패 시 None.
    """
    out_path = (mmd_path[:-4] if mmd_path.endswith(".mmd") else mmd_path) + "." + ext
    cmd = [mmdc_path, "-i", mmd_path, "-o", out_path, "-b", "transparent"]
    if ext == "png":
        cmd += ["-w", str(width)]
    try:
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=timeout, check=False,
        )
    except subprocess.TimeoutExpired:
        logger.warning("mmdc timeout (%ds) for %s.%s", timeout, mmd_path, ext)
        return None
    except Exception as e:
        logger.warning("mmdc 실행 오류 %s.%s: %s", mmd_path, ext, e)
        return None
    if result.returncode == 0 and os.path.isfile(out_path):
        return out_path
    logger.warning("mmdc 변환 실패 %s → %s: rc=%s stderr=%s",
                   mmd_path, ext, result.returncode, (result.stderr or "")[:200])
    return None


def _render_mmd_to_svg(mmdc_path: str, mmd_path: str,
                        timeout: int = 30) -> Optional[str]:
    return _render_mmd_to_format(mmdc_path, mmd_path, "svg", timeout)


def _render_mmd_to_png(mmdc_path: str, mmd_path: str,
                        timeout: int = 30) -> Optional[str]:
    return _render_mmd_to_format(mmdc_path, mmd_path, "png", timeout)


def export_flowchart_pptx(layouts: Dict[str, ScreenLayout],
                           out_path: str) -> Optional[str]:
    """모든 화면의 flowchart 를 1 PPTX 로 묶기 — 슬라이드당 1 화면.

    파이프라인: mermaid → mmdc → SVG + PNG → python-pptx 슬라이드 임베드.
    PowerPoint 의 SVG/PNG 듀얼 임베드 패턴 (svgBlip extension) 으로 vector
    edit 가능 + raster fallback. 사용자가 슬라이드에서 SVG 우클릭 →
    "도형으로 변환" 으로 mermaid 노드/엣지가 편집 가능한 PPT 도형으로.

    의존성:
      - ``python-pptx`` (PyPI). 미설치 시 None 반환 + 안내 로그.
      - ``mmdc`` (mermaid-cli, PATH). 미설치 시 None 반환.

    반환: 저장된 pptx 경로 (성공) / None (의존성 / 변환 실패).
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        logger.warning(
            "python-pptx 미설치 — pptx export skip. "
            "폐쇄망: `python -m pip download python-pptx -d .\\wheels "
            "--platform win_amd64 --python-version 311 --only-binary=:all:` "
            "후 `python -m pip install --no-index --find-links=.\\wheels python-pptx`")
        return None
    mmdc_path = _find_mmdc_executable()
    if not mmdc_path:
        logger.warning(
            "mmdc 미설치 — pptx export 위한 SVG 변환 불가. "
            "사내망에서 npm 가능하면 `npm install -g @mermaid-js/mermaid-cli`, "
            "안 되면 사내 npm registry 통해 사전 wheel/tarball 설치.")
        return None
    relevant = {rel: l for rel, l in layouts.items()
                if (l.flowchart_mermaid or "").strip()}
    if not relevant:
        logger.info("flowchart 있는 화면 0 — pptx export skip")
        return None

    import tempfile
    n_ok = 0
    n_fail = 0
    os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="flowchart_pptx_") as tmpdir:
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        for rel, layout in sorted(relevant.items()):
            code = _sanitize_mermaid_flowchart(layout.flowchart_mermaid or "")
            if not code:
                continue
            safe = re.sub(r"[^A-Za-z0-9_.-]+", "_", rel.replace("/", "__"))
            mmd_path = os.path.join(tmpdir, safe + ".mmd")
            with open(mmd_path, "w", encoding="utf-8") as f:
                f.write(code + "\n")
            png_path = _render_mmd_to_png(mmdc_path, mmd_path)
            svg_path = _render_mmd_to_svg(mmdc_path, mmd_path)
            if not png_path and not svg_path:
                n_fail += 1
                continue
            slide = prs.slides.add_slide(blank_layout)
            # 슬라이드 제목 (file_rel 기반)
            tbox = slide.shapes.add_textbox(
                Inches(0.3), Inches(0.15),
                slide_w - Inches(0.6), Inches(0.5),
            )
            tf = tbox.text_frame
            tf.text = layout.page_title or rel
            tf.paragraphs[0].runs[0].font.size = Pt(18)
            tf.paragraphs[0].runs[0].font.bold = True
            # 파일 경로 sub-text
            sub = slide.shapes.add_textbox(
                Inches(0.3), Inches(0.6),
                slide_w - Inches(0.6), Inches(0.3),
            )
            sub.text_frame.text = rel
            sub.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
            # 이미지 영역
            img_left = Inches(0.3)
            img_top = Inches(1.0)
            img_w = slide_w - Inches(0.6)
            img_h = slide_h - Inches(1.3)
            if png_path:
                pic = slide.shapes.add_picture(
                    png_path, img_left, img_top, img_w, img_h,
                )
                if svg_path:
                    _attach_svg_to_picture(pic, slide.part, svg_path)
            elif svg_path:
                # PNG 변환 실패하면 SVG 만 raw embed — 일부 PPT 버전은 native
                # 인식 못 함. fallback path 라 큰 위험은 없음.
                _add_svg_only_picture(slide, slide.part, svg_path,
                                       img_left, img_top, img_w, img_h)
            n_ok += 1
        if not n_ok:
            return None
        prs.save(out_path)
    print(f"  flowchart pptx: {n_ok} 슬라이드 → {out_path}"
          + (f" ({n_fail} 화면 변환 실패)" if n_fail else ""))
    return out_path


def _attach_svg_to_picture(pic, slide_part, svg_path: str) -> None:
    """python-pptx Picture 의 PNG blip 옆에 SVG blip extension 추가.

    ECMA-376 + Office svgBlip 패턴: ``<a:blip r:embed="rId_PNG">`` 안에
    ``<a:extLst>/<a:ext uri="{96DAC541-...}">/<asvg:svgBlip r:embed="rId_SVG"/>``.
    PowerPoint 2019+/365 가 SVG vector 로 렌더링 + 우클릭 "도형으로 변환"
    지원. 미지원 버전은 PNG fallback 자동.
    """
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.oxml.ns import qn
        from lxml import etree
    except Exception as e:
        logger.warning("svgBlip 임베드 skip (의존성): %s", e)
        return
    try:
        with open(svg_path, "rb") as f:
            svg_bytes = f.read()
    except Exception as e:
        logger.warning("SVG 읽기 실패 %s: %s", svg_path, e)
        return
    # image part 추가 + relationship 생성
    package = slide_part.package
    try:
        from pptx.parts.image import ImagePart
    except Exception:
        return
    try:
        partname = package.next_partname("/ppt/media/image%d.svg")
        image_part = ImagePart(partname, "image/svg+xml", svg_bytes, package)
        package._add_part(image_part)
    except Exception as e:
        logger.warning("SVG image part 생성 실패: %s", e)
        return
    rId_svg = slide_part.relate_to(image_part, RT.IMAGE)
    # blip XML 수정
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    asvg_ns = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    blip = pic._element.find(".//" + qn("a:blip"))
    if blip is None:
        return
    ext_lst = blip.find(qn("a:extLst"))
    if ext_lst is None:
        ext_lst = etree.SubElement(blip, qn("a:extLst"))
    ext = etree.SubElement(ext_lst, qn("a:ext"))
    ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
    svg_blip = etree.SubElement(
        ext, "{%s}svgBlip" % asvg_ns,
        nsmap={"asvg": asvg_ns},
    )
    svg_blip.set("{%s}embed" % r_ns, rId_svg)


def _add_svg_only_picture(slide, slide_part, svg_path, left, top, w, h):
    """PNG 변환 실패 시 SVG 만 image part 로 추가. PPT 버전 따라 미렌더링
    위험 있으나 fallback path 라 시도만.
    """
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.image import ImagePart
        from pptx.oxml.ns import qn
        from lxml import etree
    except Exception:
        return
    with open(svg_path, "rb") as f:
        svg_bytes = f.read()
    package = slide_part.package
    partname = package.next_partname("/ppt/media/image%d.svg")
    image_part = ImagePart(partname, "image/svg+xml", svg_bytes, package)
    package._add_part(image_part)
    rId = slide_part.relate_to(image_part, RT.IMAGE)
    # 단순 picture XML — PNG 없이 SVG 만 embed
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    pic_xml = etree.fromstring(
        f'<p:pic xmlns:p="{p_ns}" xmlns:a="{a_ns}" xmlns:r="{r_ns}">'
        f'<p:nvPicPr><p:cNvPr id="0" name="svgonly"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="{int(left)}" y="{int(top)}"/>'
        f'<a:ext cx="{int(w)}" cy="{int(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
    )
    slide.shapes._spTree.append(pic_xml)


def write_screen_html_files(out_dir: str,
                             layouts: Dict[str, ScreenLayout]) -> Dict[str, str]:
    """``{file_rel: html_path}`` 반환. 화면별 .html 파일 저장.

    레이아웃: ``out_dir/<safe_filename>.html`` (평탄). file_rel 전체를
    ``__`` 로 join 한 뒤 sanitize. timestamp dir 자체에 reponame 이
    prefix 되므로 file_rel 첫 segment (예: ``src``) 를 별도 폴더로 더
    내리지 않는다.
    """
    os.makedirs(out_dir, exist_ok=True)
    written: Dict[str, str] = {}
    for rel, layout in layouts.items():
        norm = rel.replace("\\", "/")
        parts = [p for p in norm.split("/") if p]
        joined = "__".join(parts) if parts else "screen"
        safe = re.sub(r"[^A-Za-z0-9_.-]+", "_", joined)
        path = os.path.join(out_dir, safe + ".html")
        try:
            with open(path, "w", encoding="utf-8") as f:
                f.write(render_screen_html(layout))
            written[rel] = path
        except Exception as e:
            logger.warning("screen html write 실패 %s: %s", path, e)
    return written


# ── 옵션 G stub: Playwright 스크린샷 ──────────────────────────────────


def render_screenshots_via_playwright(
    layouts: Dict[str, ScreenLayout],
    *,
    base_url: Optional[str] = None,
    out_dir: str = "output/legacy_analysis/screenshots",
) -> Dict[str, str]:
    """별도 옵트인 (`--render-screenshots`). 사용자 PC 에 React 빌드/실행
    + Playwright 설치된 환경에서만 동작.

    현재는 스텁 — 실제 구현은 follow-up. 호출 시 안내 메시지 emit.
    """
    print("  screenshot rendering: --render-screenshots 옵션은 현재 stub.")
    print("  실제 구현은 사용자 PC 에 React 빌드 + Playwright 셋업 필요.")
    print("  follow-up PR 에서 base_url 기반 자동 스크린샷 추가 예정.")
    return {}
